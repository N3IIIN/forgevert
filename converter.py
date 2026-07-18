#!/usr/bin/env python3
"""
Universal File Converter — Konvertierungs-Engine
~300+ Formate, BFS-Pfadsuche, externe Prozesse (ffmpeg, Blender, LibreOffice, Pandoc)
"""
import os, sys, base64, tempfile, shutil, subprocess, json
from pathlib import Path
from collections import deque

# ── Optionale Bibliotheken erkennen ──────────────────────────────────────────

def _try_import(name):
    try: __import__(name); return True
    except Exception: return False

_PIL        = _try_import("PIL")
_TRIMESH    = _try_import("trimesh")
_PANDAS     = _try_import("pandas")
_PYPDF      = _try_import("pypdf")
_DOCX       = _try_import("docx")
_REPORTLAB  = _try_import("reportlab")
_EZDXF      = _try_import("ezdxf")
_MARKDOWN   = _try_import("markdown")
_PPTX       = _try_import("pptx")
_ODFPY      = _try_import("odf")
_RHINO3DM   = _try_import("rhino3dm")
_IFCOS      = _try_import("ifcopenshell")
_PYEMB      = _try_import("pyembroidery")
_MIDO       = _try_import("mido")
_FONTTOOLS  = _try_import("fontTools")
_EBOOKLIB   = _try_import("ebooklib")
_NBFORMAT   = _try_import("nbformat")
_NBCONVERT  = _try_import("nbconvert")
_PYARROW    = _try_import("pyarrow")
_PYSHP      = _try_import("shapefile")
_GPXPY      = _try_import("gpxpy")
_FASTKML    = _try_import("fastkml")
_LASPY      = _try_import("laspy")
_PYDICOM    = _try_import("pydicom")
_NIBABEL    = _try_import("nibabel")
_RAWPY      = _try_import("rawpy")
_PILHEIF    = _try_import("pillow_heif")
_GERBER     = _try_import("gerber")
_PY7ZR      = _try_import("py7zr")
_RARFILE    = _try_import("rarfile")
_VSDX       = _try_import("vsdx")
_EXTRACTMSG = _try_import("extract_msg")
_PYMUPDF    = _try_import("fitz")
_REMBG      = _try_import("rembg")
_XLWT       = _try_import("xlwt")
_FIONA      = _try_import("fiona")
_MUSIC21    = _try_import("music21")
_H5PY       = _try_import("h5py")
_NETCDF4    = _try_import("netCDF4")
_RDKIT      = _try_import("rdkit")
_YAML       = _try_import("yaml")
_ICALENDAR  = _try_import("icalendar")
_QRCODE     = _try_import("qrcode")
_PDF2DOCX   = _try_import("pdf2docx")

def _has_cmd(*cmds):
    for cmd in cmds:
        try:
            subprocess.run([cmd, "--version"], capture_output=True, timeout=5)
            return cmd
        except Exception:
            try:
                subprocess.run([cmd, "-version"], capture_output=True, timeout=5)
                return cmd
            except Exception:
                continue
    return None

_FFMPEG    = bool(_has_cmd("ffmpeg"))

# ── Vorab-Laden schwerer Bibliotheken (einmalig beim Start) ──────────────────
# Synchron geladen, bevor der Server Anfragen annimmt → kein GIL-Konflikt
if _PIL:
    try: from PIL import Image as _PIL_Image, ImageOps, ImageFilter  # noqa: F401
    except Exception: pass
if _TRIMESH:
    try: import trimesh as _trimesh_mod, numpy as _np_mod  # noqa: F401
    except Exception: pass
if _EZDXF:
    try: import ezdxf as _ezdxf_mod  # noqa: F401
    except Exception: pass
if _PYMUPDF:
    try: import fitz as _fitz_mod  # noqa: F401
    except Exception: pass
if _PANDAS:
    try: import pandas as _pd_mod  # noqa: F401
    except Exception: pass

def _find_soffice():
    cmd = _has_cmd("libreoffice", "soffice")
    if cmd: return cmd
    import platform
    if platform.system() == "Windows":
        for p in [
            r"C:\Program Files\LibreOffice\program\soffice.exe",
            r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        ]:
            if Path(p).exists(): return p
    return None

_SOFFICE = _find_soffice()

def _find_lo_python():
    """Python mit UNO-Modul finden (Windows: LibreOffice-Bundle, Linux: python3-uno)"""
    import platform
    if platform.system() == "Windows":
        for p in [
            r"C:\Program Files\LibreOffice\program\python.exe",
            r"C:\Program Files (x86)\LibreOffice\program\python.exe",
        ]:
            if Path(p).exists():
                return p
        return None
    # Linux/Docker: system python3 mit python3-uno prüfen
    for p in ["/usr/bin/python3", "/usr/local/bin/python3"]:
        if Path(p).exists():
            try:
                r = subprocess.run([p, "-c", "import uno"],
                                   capture_output=True, timeout=5)
                if r.returncode == 0:
                    return p
            except Exception:
                pass
    return None

_LO_PYTHON = _find_lo_python()
_LO_PORT   = 2002
_LO_LISTENER_PROC  = None
_LO_LISTENER_LOCK  = None  # wird nach threading-Import gesetzt

def _lo_listener_alive():
    import socket as _sock
    try:
        s = _sock.create_connection(("127.0.0.1", _LO_PORT), timeout=1.0)
        s.close(); return True
    except Exception:
        return False

def _ensure_lo_listener():
    """Startet LibreOffice-Listener falls nicht aktiv. Thread-sicher."""
    global _LO_LISTENER_PROC, _LO_LISTENER_LOCK
    if not _SOFFICE or not _LO_PYTHON:
        return False
    if _lo_listener_alive():
        return True
    if _LO_LISTENER_LOCK is None:
        import threading as _t
        _LO_LISTENER_LOCK = _t.Lock()
    with _LO_LISTENER_LOCK:
        if _lo_listener_alive():
            return True
        try:
            _LO_LISTENER_PROC = subprocess.Popen([
                _SOFFICE, "--headless", "--norestore", "--nofirststartwizard",
                f"--accept=socket,host=127.0.0.1,port={_LO_PORT};urp;StarOffice.ServiceManager",
            ])
            import time as _t2
            for _ in range(24):          # max 12s warten
                _t2.sleep(0.5)
                if _lo_listener_alive():
                    return True
        except Exception:
            pass
        return False
_BLENDER   = bool(_has_cmd("blender"))
_PANDOC    = bool(_has_cmd("pandoc"))
_INKSCAPE  = bool(_has_cmd("inkscape"))
_GHOSTSCR  = bool(_has_cmd("gs", "gswin64c"))

# ── Format-Normalisierung ─────────────────────────────────────────────────────

_ALIASES = {
    # Bild
    "jpg":"jpeg", "jpe":"jpeg", "jfif":"jpeg",
    "tif":"tiff",
    "jxl":"jpeg_xl",
    # CAD
    "stp":"step", "p21":"step",
    "igs":"iges",
    "gltf":"glb",
    "wrl":"vrml",
    "x_t":"parasolid", "x_b":"parasolid", "xmt_txt":"parasolid",
    # Audio
    "m4a":"aac", "m4b":"aac",
    # Video
    "mpeg":"mpg", "m2v":"mpg", "m4v":"mp4",
    # Dokument
    "yml":"yaml", "htm":"html", "tgz":"tar_gz", "tbz2":"tar_bz2", "txz":"tar_xz",
    "md":"markdown",
    "markdown":"markdown",
    # GIS
    "geojson":"geojson",
    # Stickerei
    "jef+":"jef_plus",
    # Font
    "ttf":"ttf", "otf":"otf",
}

def norm(fmt: str) -> str:
    f = fmt.lower().lstrip(".")
    return _ALIASES.get(f, f)

def ext_of(path: str) -> str:
    name = Path(path).name.lower()
    if name.endswith(".tar.gz"):  return "tar_gz"
    if name.endswith(".tar.bz2"): return "tar_bz2"
    if name.endswith(".tar.xz"):  return "tar_xz"
    return norm(Path(path).suffix)

def _b64(path: str) -> str:
    with open(path, "rb") as f:
        return base64.b64encode(f.read()).decode()

def _run(cmd: list, timeout=300, env=None) -> subprocess.CompletedProcess:
    return subprocess.run(cmd, capture_output=True, timeout=timeout, env=env)

# ── Conversion Functions ──────────────────────────────────────────────────────

# ══════════ BILDER ══════════

def _img_to_img(src, dst, **kw):
    from PIL import Image
    img = Image.open(src)
    ext = norm(Path(dst).suffix).upper()

    # Resize
    rw = kw.get("resize_w")
    rh = kw.get("resize_h")
    if rw or rh:
        orig_w, orig_h = img.size
        if rw and rh:
            img = img.resize((int(rw), int(rh)), Image.LANCZOS)
        elif rw:
            ratio = int(rw) / orig_w
            img = img.resize((int(rw), int(orig_h * ratio)), Image.LANCZOS)
        else:
            ratio = int(rh) / orig_h
            img = img.resize((int(orig_w * ratio), int(rh)), Image.LANCZOS)

    if ext == "JPEG":
        img = img.convert("RGB")
    elif ext in ("BMP", "PCX"):
        img = img.convert("RGB")
    elif ext in ("PPM", "PBM"):
        img = img.convert("RGB")
    elif ext == "PGM":
        img = img.convert("L")
    elif ext == "ICO":
        img.save(dst, format="ICO", sizes=[(256,256),(128,128),(64,64),(32,32),(16,16)])
        return

    quality = kw.get("quality")
    save_kw = {}
    if quality is not None:
        q = max(1, min(95, int(quality)))
        if ext in ("JPEG",):
            save_kw["quality"] = q
            save_kw["optimize"] = True
        elif ext == "WEBP":
            save_kw["quality"] = q
        elif ext == "PNG":
            # PNG hat kein quality, aber compress_level 0-9
            save_kw["compress_level"] = max(0, min(9, 9 - q // 11))
    img.save(dst, **save_kw)

def _img_to_pdf(src, dst, **kw):
    from PIL import Image
    Image.open(src).convert("RGB").save(dst, "PDF", resolution=150)

def _pdf_to_img(src, dst, **kw):
    ext = norm(Path(dst).suffix)
    if _PYMUPDF:
        import fitz
        doc = fitz.open(src)
        pix = doc[0].get_pixmap(matrix=fitz.Matrix(2, 2))
        pix.save(dst)
        return
    try:
        import pypdfium2 as pdfium
        pdf = pdfium.PdfDocument(src)
        pil = pdf[0].render(scale=2).to_pil()
        if ext == "jpeg": pil = pil.convert("RGB")
        pil.save(dst)
        return
    except ImportError:
        pass
    from PIL import Image
    img = Image.open(src)
    if ext == "jpeg": img = img.convert("RGB")
    img.save(dst)

def _raw_to_img(src, dst, **kw):
    """Kamera-RAW (CR3/CR2/NEF/ARW/RAF/DNG/ORF/RW2...) → Bild via rawpy"""
    import rawpy
    from PIL import Image
    with rawpy.imread(src) as raw:
        rgb = raw.postprocess(use_camera_wb=True, output_bps=8)
    ext = norm(Path(dst).suffix).upper()
    img = Image.fromarray(rgb)
    if ext == "JPEG": img = img.convert("RGB")
    img.save(dst)

def _heic_to_img(src, dst, **kw):
    """HEIC/AVIF → Bild via pillow-heif"""
    from pillow_heif import register_heif_opener
    from PIL import Image
    register_heif_opener()
    img = Image.open(src)
    ext = norm(Path(dst).suffix).upper()
    if ext == "JPEG": img = img.convert("RGB")
    img.save(dst)

def _eps_to_png(src, dst, **kw):
    """EPS → PNG via Ghostscript"""
    r = _run(["gs", "-dNOPAUSE", "-dBATCH", "-dSAFER",
               "-sDEVICE=png16m", "-r150", f"-sOutputFile={dst}", src])
    if r.returncode != 0:
        raise RuntimeError(f"Ghostscript: {r.stderr.decode()[-300:]}")

def _svg_to_png(src, dst, **kw):
    """SVG → PNG via Inkscape oder cairosvg"""
    if _INKSCAPE:
        r = _run(["inkscape", "--export-type=png", f"--export-filename={dst}", src])
        if r.returncode == 0: return
    try:
        import cairosvg
        cairosvg.svg2png(url=src, write_to=dst)
    except ImportError:
        raise RuntimeError("SVG→PNG: pip install cairosvg  oder  inkscape installieren")

def _psd_to_png(src, dst, **kw):
    """PSD/PSB → PNG via psd-tools"""
    from psd_tools import PSDImage
    img = PSDImage.open(src)
    img.compose().save(dst)

def _img_to_txt_ocr(src, dst, **kw):
    """Bild → TXT via Tesseract OCR"""
    import pytesseract
    from PIL import Image
    import platform
    if _PILHEIF:
        try:
            import pillow_heif
            pillow_heif.register_heif_opener()
        except Exception:
            pass
    if platform.system() == "Windows":
        for tpath in [r"C:\Program Files\Tesseract-OCR\tesseract.exe",
                      r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe"]:
            if Path(tpath).exists():
                pytesseract.pytesseract.tesseract_cmd = tpath
                break
    img = Image.open(src)
    # Tesseract akzeptiert kein HEIF-Format-Tag → Bild neu erstellen (strips format attribute)
    import numpy as np
    img = Image.fromarray(np.array(img.convert("RGB")))
    text = pytesseract.image_to_string(img, lang="deu+eng")
    if not text.strip():
        text = pytesseract.image_to_string(img, lang="eng")
    Path(dst).write_text(text.strip() or "(kein Text erkannt)", encoding="utf-8")

# ══════════ DOKUMENTE ══════════

def _libreoffice(src, dst, **kw):
    """Office-Format → Zielformat via LibreOffice.
    Schneller Pfad: persistenter UNO-Listener.
    Fallback: klassischer --convert-to Prozess."""
    ext = Path(dst).suffix.lstrip(".")

    # ── Schneller Pfad: UNO-Listener ─────────────────────────────────────────
    _LO_HELPER = Path(__file__).parent / "lo_helper.py"
    if _LO_PYTHON and _LO_HELPER.exists() and _ensure_lo_listener():
        r = _run([_LO_PYTHON, str(_LO_HELPER),
                  str(Path(src).resolve()), str(Path(dst).resolve()),
                  "127.0.0.1", str(_LO_PORT)], timeout=60)
        if Path(dst).exists():
            return
        # UNO schlug fehl → weiter zum Fallback

    # ── Fallback: neuer Prozess pro Konvertierung ─────────────────────────────
    out_dir = str(Path(dst).parent)
    r = _run([_SOFFICE, "--headless", "--norestore", "--nofirststartwizard",
              "--convert-to", ext, "--outdir", out_dir, src], timeout=300)
    expected = Path(out_dir) / (Path(src).stem + "." + ext)
    if expected.exists() and str(expected) != str(dst):
        shutil.move(str(expected), dst)
    if not Path(dst).exists():
        raise RuntimeError(f"LibreOffice Fehler: {r.stderr.decode(errors='replace')[:400]}")

def _pandoc(src, dst, **kw):
    """Markdown/RST/HTML/DOCX/EPUB ↔ alle via Pandoc"""
    r = _run(["pandoc", src, "-o", dst, "--standalone"], timeout=120)
    if r.returncode != 0:
        raise RuntimeError(f"Pandoc: {r.stderr.decode(errors='replace')[-300:]}")

def _docx_to_txt(src, dst, **kw):
    import docx
    text = "\n".join(p.text for p in docx.Document(src).paragraphs)
    Path(dst).write_text(text, encoding="utf-8")

def _docx_to_html(src, dst, **kw):
    import docx
    doc = docx.Document(src)
    parts = ["<!DOCTYPE html><html><body>"]
    for p in doc.paragraphs:
        if p.style and p.style.name.startswith("Heading"):
            lvl = p.style.name.split()[-1]
            lvl = lvl if lvl.isdigit() else "1"
            parts.append(f"<h{lvl}>{p.text}</h{lvl}>")
        elif p.text.strip():
            parts.append(f"<p>{p.text}</p>")
    parts.append("</body></html>")
    Path(dst).write_text("\n".join(parts), encoding="utf-8")

def _txt_to_docx(src, dst, **kw):
    import docx
    doc = docx.Document()
    for line in Path(src).read_text(encoding="utf-8", errors="replace").split("\n"):
        doc.add_paragraph(line)
    doc.save(dst)

def _docx_to_pdf_native(src, dst, **kw):
    """python-docx + reportlab fallback (kein LibreOffice/LaTeX nötig)."""
    import docx
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import cm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer

    doc = docx.Document(src)
    pdf = SimpleDocTemplate(dst, pagesize=A4,
                            leftMargin=2.5*cm, rightMargin=2.5*cm,
                            topMargin=2.5*cm, bottomMargin=2.5*cm)
    styles = getSampleStyleSheet()
    story = []

    for p in doc.paragraphs:
        text = p.text.strip()
        if not text:
            story.append(Spacer(1, 6))
            continue
        style_name = "Normal"
        if p.style and p.style.name.startswith("Heading"):
            lvl = p.style.name.split()[-1]
            if lvl.isdigit() and 1 <= int(lvl) <= 3:
                style_name = f"Heading{lvl}"
        import html as _html_mod
        safe = _html_mod.escape(text)
        if p.style and "Bold" in (p.style.name or ""):
            safe = f"<b>{safe}</b>"
        story.append(Paragraph(safe, styles[style_name]))
        story.append(Spacer(1, 3))

    if not story:
        story.append(Paragraph("(Leer)", styles["Normal"]))
    pdf.build(story)

def _pdf_to_docx(src, dst, **kw):
    from pdf2docx import Converter as PdfConverter
    cv = PdfConverter(src)
    cv.convert(dst, start=0, end=None)
    cv.close()

def _txt_to_pdf(src, dst, **kw):
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import A4
    c = canvas.Canvas(dst, pagesize=A4)
    W, H = A4
    y = H - 50
    c.setFont("Helvetica", 10)
    for line in Path(src).read_text(encoding="utf-8", errors="replace").split("\n"):
        if y < 50: c.showPage(); y = H - 50; c.setFont("Helvetica", 10)
        c.drawString(40, y, line[:130])
        y -= 13
    c.save()

def _pdf_to_txt(src, dst, **kw):
    if _PYMUPDF:
        import fitz
        doc = fitz.open(src)
        text = "\n\n".join(page.get_text() for page in doc)
    elif _PYPDF:
        import pypdf
        reader = pypdf.PdfReader(src)
        text = "\n\n".join(p.extract_text() or "" for p in reader.pages)
    else:
        raise RuntimeError("pip install pymupdf  oder  pypdf")
    Path(dst).write_text(text, encoding="utf-8")

def _md_to_html(src, dst, **kw):
    text = Path(src).read_text(encoding="utf-8")
    if _MARKDOWN:
        import markdown
        body = markdown.markdown(text, extensions=["tables","fenced_code","nl2br"])
    else:
        body = text.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;").replace("\n","<br>")
    Path(dst).write_text(f"<!DOCTYPE html><html><head><meta charset='utf-8'></head><body>{body}</body></html>",
                         encoding="utf-8")

def _html_to_pdf(src, dst, **kw):
    try:
        import weasyprint
        weasyprint.HTML(filename=src).write_pdf(dst); return
    except ImportError: pass
    if _PANDOC:
        _pandoc(src, dst); return
    raise RuntimeError("HTML→PDF: pip install weasyprint")

def _txt_to_html(src, dst, **kw):
    escaped = Path(src).read_text(encoding="utf-8",errors="replace").replace("&","&amp;").replace("<","&lt;").replace(">","&gt;")
    Path(dst).write_text(f"<!DOCTYPE html><html><body><pre>{escaped}</pre></body></html>", encoding="utf-8")

def _html_to_txt(src, dst, **kw):
    import re
    html = Path(src).read_text(encoding="utf-8", errors="replace")
    Path(dst).write_text(re.sub(r"<[^>]+>", "", html), encoding="utf-8")

def _epub_to_txt(src, dst, **kw):
    import ebooklib
    from ebooklib import epub
    import re
    book = epub.read_epub(src)
    texts = []
    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        text = re.sub(r"<[^>]+>", "", item.get_content().decode("utf-8", errors="replace"))
        texts.append(text)
    Path(dst).write_text("\n\n".join(texts), encoding="utf-8")

def _epub_to_html(src, dst, **kw):
    import ebooklib
    from ebooklib import epub
    book = epub.read_epub(src)
    parts = ["<!DOCTYPE html><html><body>"]
    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        parts.append(item.get_content().decode("utf-8", errors="replace"))
    parts.append("</body></html>")
    Path(dst).write_text("\n".join(parts), encoding="utf-8")

def _ipynb_to_html(src, dst, **kw):
    import nbformat
    from nbconvert import HTMLExporter
    nb = nbformat.read(src, as_version=4)
    body, _ = HTMLExporter().from_notebook_node(nb)
    Path(dst).write_text(body, encoding="utf-8")

def _ipynb_to_py(src, dst, **kw):
    import nbformat
    from nbconvert import PythonExporter
    nb = nbformat.read(src, as_version=4)
    body, _ = PythonExporter().from_notebook_node(nb)
    Path(dst).write_text(body, encoding="utf-8")

def _ipynb_to_md(src, dst, **kw):
    import nbformat
    from nbconvert import MarkdownExporter
    nb = nbformat.read(src, as_version=4)
    body, _ = MarkdownExporter().from_notebook_node(nb)
    Path(dst).write_text(body, encoding="utf-8")

def _pptx_to_html(src, dst, **kw):
    """PPTX → HTML (alle Folientexte)"""
    from pptx import Presentation
    import html as html_mod
    prs = Presentation(src)
    parts = ['<!DOCTYPE html><html><head><meta charset="utf-8">',
             '<style>.slide{border:1px solid #ccc;margin:20px auto;max-width:900px;',
             'padding:30px;page-break-after:always;min-height:200px}</style></head><body>']
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f'<div class="slide"><h2>Folie {i}</h2>')
        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text: continue
                is_bold = any(r.font.bold for r in para.runs if r.font.bold is not None)
                tag = "b" if is_bold else "p"
                parts.append(f'<{tag}>{html_mod.escape(text)}</{tag}>')
        parts.append('</div>')
    parts.append('</body></html>')
    Path(dst).write_text(''.join(parts), encoding='utf-8')

def _pptx_to_txt(src, dst, **kw):
    """PPTX → TXT (alle Folientexte extrahiert)"""
    from pptx import Presentation
    prs = Presentation(src)
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"=== Folie {i} ===")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t: lines.append(t)
        lines.append("")
    Path(dst).write_text("\n".join(lines), encoding="utf-8")

def _pptx_to_pdf(src, dst, **kw):
    if _SOFFICE: _libreoffice(src, dst); return
    # Pure-Python-Fallback: PPTX → HTML → PDF
    html_tmp = src + "_pptx_tmp.html"
    try:
        _pptx_to_html(src, html_tmp)
        _html_to_pdf(html_tmp, dst)
    finally:
        try: os.unlink(html_tmp)
        except: pass

# ══════════ TABELLEN / DATEN ══════════

_DATA_READERS = None  # cached on first use
_DATA_WRITERS = None  # cached on first use

def _pd_read_json(p):
    """JSON → DataFrame; normalisiert flache Dicts und einfache Objekte."""
    import json
    import pandas as pd
    raw = Path(p).read_text(encoding="utf-8", errors="replace")
    data = json.loads(raw)
    if isinstance(data, list):
        return pd.DataFrame(data)
    if isinstance(data, dict):
        # Alle Werte skalare → Einzelzeile
        if all(not isinstance(v, (dict, list)) for v in data.values()):
            return pd.DataFrame([data])
        # Ein Schlüssel, dessen Wert eine Liste von Dicts ist → unwrap
        vals = list(data.values())
        if len(data) == 1 and isinstance(vals[0], list):
            return pd.DataFrame(vals[0])
        # Spalten-orientiertes Dict {col: [v1,v2,...]}
        first = vals[0] if vals else None
        if isinstance(first, list):
            return pd.DataFrame(data)
        # Fallback: jeder Wert als Zeile (z.B. DXF entities-Dict)
        return pd.DataFrame([{k: str(v) if isinstance(v, (dict, list)) else v
                               for k, v in data.items()}])
    return pd.read_json(p)

def _pd_read_toml(p):
    """TOML → DataFrame; unwrapped array-of-tables oder flat dict."""
    import pandas as pd
    try:
        import tomllib
    except ImportError:
        import tomli as tomllib
    data = tomllib.loads(Path(p).read_text(encoding="utf-8"))
    # Array-of-tables: {"data": [{...},{...}]} → unwrap
    vals = list(data.values())
    if len(data) == 1 and isinstance(vals[0], list):
        return pd.DataFrame(vals[0])
    # Flaches Dict → Einzelzeile
    if all(not isinstance(v, (dict, list)) for v in data.values()):
        return pd.DataFrame([data])
    # Allgemeiner Fallback: stringify nested values
    return pd.DataFrame([{k: str(v) if isinstance(v, (dict, list)) else v
                          for k, v in data.items()}])

def _get_data_readers():
    import pandas as pd
    readers = {
        "csv":     lambda p: pd.read_csv(p),
        "tsv":     lambda p: pd.read_csv(p, sep="\t"),
        "json":    _pd_read_json,
        "xlsx":    lambda p: pd.read_excel(p, engine="openpyxl"),
        "xls":     lambda p: pd.read_excel(p),
        "xml":     lambda p: pd.read_xml(p),
        "parquet": lambda p: pd.read_parquet(p),
        "feather": lambda p: pd.read_feather(p),
        "orc":     lambda p: pd.read_orc(p),
        "hdf5":    lambda p: pd.read_hdf(p),
        "toml":    _pd_read_toml,
    }
    return readers

def _write_xls(df, path):
    import xlwt
    wb = xlwt.Workbook(encoding="utf-8")
    ws = wb.add_sheet("Sheet1")
    for col_i, col_name in enumerate(df.columns):
        ws.write(0, col_i, str(col_name))
    for row_i, row in enumerate(df.itertuples(index=False), 1):
        for col_i, val in enumerate(row):
            if hasattr(val, 'item'):
                val = val.item()
            if not isinstance(val, (int, float, str, bool, type(None))):
                val = str(val)
            ws.write(row_i, col_i, val)
    wb.save(str(path))

def _get_data_writers():
    import pandas as pd
    writers = {
        "csv":      lambda df, p: df.to_csv(p, index=False),
        "tsv":      lambda df, p: df.to_csv(p, sep="\t", index=False),
        "json":     lambda df, p: df.to_json(p, orient="records", indent=2, force_ascii=False),
        "xlsx":     lambda df, p: df.to_excel(p, index=False, engine="openpyxl"),
        "xml":      lambda df, p: df.to_xml(p, index=False),
        "parquet":  lambda df, p: df.to_parquet(p, index=False),
        "feather":  lambda df, p: df.to_feather(p),
        "html":     lambda df, p: Path(p).write_text(
            f"<!DOCTYPE html><html><body>{df.to_html(index=False)}</body></html>", encoding="utf-8"),
        "markdown": lambda df, p: Path(p).write_text(df.to_markdown(index=False) or "*keine Daten*", encoding="utf-8"),
        "txt":      lambda df, p: Path(p).write_text(df.to_string(index=False), encoding="utf-8"),
    }
    if _XLWT:
        writers["xls"] = _write_xls
    if _PYARROW:
        try:
            import pyarrow.orc as orc_mod
            writers["orc"] = lambda df, p: df.to_orc(p, index=False)
        except Exception:
            pass
    try:
        import tomli_w
        def _write_toml_df(df, path):
            records = []
            for row in df.where(df.notna(), None).to_dict(orient="records"):
                records.append({k: (v.item() if hasattr(v, "item") else
                                    (str(v) if not isinstance(v, (int, float, bool, str, type(None))) else v))
                                 for k, v in row.items() if v is not None})
            Path(path).write_text(tomli_w.dumps({"data": records}) or "# keine Daten\n", encoding="utf-8")
        writers["toml"] = _write_toml_df
    except ImportError:
        pass
    return writers

def _data_convert(src, dst, **kw):
    global _DATA_READERS, _DATA_WRITERS
    if _DATA_READERS is None: _DATA_READERS = _get_data_readers()
    if _DATA_WRITERS is None: _DATA_WRITERS = _get_data_writers()
    src_ext = norm(Path(src).suffix)
    dst_ext = norm(Path(dst).suffix)
    if src_ext not in _DATA_READERS: raise ValueError(f"Unbekanntes Quellformat: {src_ext}")
    if dst_ext not in _DATA_WRITERS: raise ValueError(f"Unbekanntes Zielformat: {dst_ext}")
    _DATA_WRITERS[dst_ext](_DATA_READERS[src_ext](src), dst)

def _sqlite_to_json(src, dst, **kw):
    import sqlite3, json
    conn = sqlite3.connect(src)
    result = {}
    for table in conn.execute("SELECT name FROM sqlite_master WHERE type='table'").fetchall():
        t = table[0]
        rows = conn.execute(f"SELECT * FROM [{t}]").fetchall()
        cols = [d[0] for d in conn.execute(f"SELECT * FROM [{t}] LIMIT 0").description or []]
        result[t] = [dict(zip(cols, row)) for row in rows]
    conn.close()
    Path(dst).write_text(json.dumps(result, indent=2, default=str), encoding="utf-8")

def _sqlite_to_csv(src, dst, **kw):
    import sqlite3, csv
    conn = sqlite3.connect(src)
    tables = conn.execute("SELECT name FROM sqlite_master WHERE type='table'").fetchall()
    if not tables: raise ValueError("SQLite-Datenbank ist leer")
    table = tables[0][0]
    cur = conn.execute(f"SELECT * FROM [{table}]")
    cols = [d[0] for d in cur.description]
    with open(dst, "w", newline="", encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(cols)
        writer.writerows(cur.fetchall())
    conn.close()

# ══════════ AUDIO (ffmpeg) ══════════

def _audio_convert(src, dst, **kw):
    if not _FFMPEG: raise RuntimeError("ffmpeg nicht installiert")
    r = _run(["ffmpeg", "-y", "-i", src, dst], timeout=600)
    if r.returncode != 0: raise RuntimeError(f"ffmpeg: {r.stderr.decode(errors='replace')[-400:]}")

def _midi_to_wav(src, dst, **kw):
    """MIDI → WAV via FluidSynth"""
    sf2_paths = ["/usr/share/sounds/sf2/FluidR3_GM.sf2",
                 "/usr/share/soundfonts/FreePats-GeneralUser.sf2"]
    sf2 = next((p for p in sf2_paths if Path(p).exists()), None)
    if not sf2:
        raise RuntimeError("FluidSynth SoundFont nicht gefunden: apt install fluid-soundfont-gm")
    r = _run(["fluidsynth", "-ni", sf2, src, "-F", dst, "-r", "44100"], timeout=120)
    if r.returncode != 0: raise RuntimeError(f"FluidSynth: {r.stderr.decode()[-300:]}")

def _midi_to_json(src, dst, **kw):
    """MIDI → JSON (Noten-Daten)"""
    import mido
    mid = mido.MidiFile(src)
    tracks = []
    for track in mid.tracks:
        msgs = []
        for msg in track:
            msgs.append(msg.dict())
        tracks.append({"name": track.name, "messages": msgs})
    Path(dst).write_text(json.dumps({
        "type": mid.type, "ticks_per_beat": mid.ticks_per_beat, "tracks": tracks
    }, indent=2), encoding="utf-8")

# ══════════ VIDEO (ffmpeg) ══════════

def _video_convert(src, dst, **kw):
    if not _FFMPEG: raise RuntimeError("ffmpeg nicht installiert")
    ext = norm(Path(dst).suffix)
    extra = []
    if ext == "gif":
        extra = ["-vf", "fps=10,scale=640:-1:flags=lanczos", "-loop", "0"]
    elif ext == "mp4":
        extra = ["-c:v", "libx264", "-crf", "23", "-preset", "fast", "-movflags", "+faststart"]
    elif ext == "webm":
        extra = ["-c:v", "libvpx-vp9", "-crf", "30", "-b:v", "0"]
    elif ext == "3gp":
        # Prüfe ob Audio-Stream vorhanden
        _probe = subprocess.run(
            ["ffprobe", "-v", "quiet", "-select_streams", "a:0",
             "-show_entries", "stream=codec_name", "-of", "csv=p=0", src],
            capture_output=True, timeout=10
        )
        has_audio = bool(_probe.stdout.strip())
        extra = ["-c:v", "libx264", "-profile:v", "baseline", "-level", "3.0",
                 "-vf", "scale=max(2\\,trunc(iw/2)*2):max(2\\,trunc(ih/2)*2)"]
        if has_audio:
            extra += ["-c:a", "aac", "-ar", "8000", "-ac", "1", "-b:a", "12k"]
        else:
            extra += ["-an"]
    elif ext == "apng":
        extra = ["-f", "apng", "-plays", "0"]
    r = _run(["ffmpeg", "-y", "-i", src] + extra + [dst], timeout=1800)
    if r.returncode != 0: raise RuntimeError(f"ffmpeg: {r.stderr.decode(errors='replace')[-400:]}")

def _video_to_audio(src, dst, **kw):
    if not _FFMPEG: raise RuntimeError("ffmpeg nicht installiert")
    probe = subprocess.run(
        ["ffprobe", "-v", "quiet", "-select_streams", "a:0",
         "-show_entries", "stream=codec_name", "-of", "csv=p=0", src],
        capture_output=True, timeout=10
    )
    if not probe.stdout.strip():
        raise RuntimeError("Keine Audio-Spur im Quellformat vorhanden")
    r = _run(["ffmpeg", "-y", "-i", src, "-vn", "-q:a", "2", dst], timeout=600)
    if r.returncode != 0: raise RuntimeError(f"ffmpeg: {r.stderr.decode(errors='replace')[-400:]}")

def _video_thumbnail(src, dst, **kw):
    """Video → Thumbnail (erster Frame)"""
    if not _FFMPEG: raise RuntimeError("ffmpeg nicht installiert")
    r = _run(["ffmpeg", "-y", "-i", src, "-vframes", "1", "-q:v", "2", dst], timeout=60)
    if r.returncode != 0: raise RuntimeError(f"ffmpeg: {r.stderr.decode(errors='replace')[-300:]}")

# ══════════ 3D / MESH ══════════

def _mesh_convert(src, dst, **kw):
    import trimesh
    m = trimesh.load(src, force="mesh")
    m.export(dst)

def _mesh_to_3mf(src, dst, **kw):
    """Mesh → 3MF (Fusion 360 / PrusaSlicer / Bambu kompatibel)"""
    import trimesh, zipfile
    m = trimesh.load(src, force="mesh")
    verts, faces = m.vertices, m.faces

    # ── 3D/3dmodel.model ────────────────────────────────────────────────────
    vert_lines = "\n      ".join(
        f'<vertex x="{v[0]:.6f}" y="{v[1]:.6f}" z="{v[2]:.6f}"/>' for v in verts
    )
    tri_lines = "\n      ".join(
        f'<triangle v1="{f[0]}" v2="{f[1]}" v3="{f[2]}"/>' for f in faces
    )
    model_xml = f"""<?xml version="1.0" encoding="UTF-8"?>
<model unit="millimeter" xml:lang="en-US"
  xmlns="http://schemas.microsoft.com/3dmanufacturing/core/2015/02">
  <metadata name="Application">Forgevert</metadata>
  <metadata name="Description">Converted by forgevert.onrender.com</metadata>
  <resources>
    <object id="1" type="model">
      <mesh>
        <vertices>
      {vert_lines}
        </vertices>
        <triangles>
      {tri_lines}
        </triangles>
      </mesh>
    </object>
  </resources>
  <build>
    <item objectid="1"/>
  </build>
</model>"""

    content_types = """<?xml version="1.0" encoding="UTF-8"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="model" ContentType="application/vnd.ms-package.3dmanufacturing-3dmodel+xml"/>
</Types>"""

    rels = """<?xml version="1.0" encoding="UTF-8"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Target="/3D/3dmodel.model" Id="rel-1"
    Type="http://schemas.microsoft.com/3dmanufacturing/2013/01/3dmodel"/>
</Relationships>"""

    with zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("[Content_Types].xml", content_types)
        zf.writestr("_rels/.rels", rels)
        zf.writestr("3D/3dmodel.model", model_xml)

def _mesh_to_amf(src, dst, **kw):
    """Mesh → AMF (Additive Manufacturing Format, XML)"""
    import trimesh
    m = trimesh.load(src, force="mesh")
    verts, faces = m.vertices, m.faces
    lines = [
        '<?xml version="1.0" encoding="utf-8"?>',
        '<amf unit="millimeter" version="1.1">',
        ' <object id="1"><mesh><vertices>',
    ]
    for v in verts:
        lines.append(f'  <vertex><coordinates>'
                     f'<x>{v[0]:.6f}</x><y>{v[1]:.6f}</y><z>{v[2]:.6f}</z>'
                     f'</coordinates></vertex>')
    lines.append(' </vertices><volume>')
    for f in faces:
        lines.append(f'  <triangle>'
                     f'<v1>{f[0]}</v1><v2>{f[1]}</v2><v3>{f[2]}</v3>'
                     f'</triangle>')
    lines.append(' </volume></mesh></object></amf>')
    Path(dst).write_text('\n'.join(lines), encoding='utf-8')

def _mesh_to_x3d(src, dst, **kw):
    """Mesh → X3D (Web3D Standard, XML)"""
    import trimesh
    m = trimesh.load(src, force="mesh")
    verts, faces = m.vertices, m.faces
    coord_str = " ".join(f"{v[0]:.4f} {v[1]:.4f} {v[2]:.4f}" for v in verts)
    idx_str   = " ".join(f"{f[0]} {f[1]} {f[2]} -1" for f in faces)
    x3d = (
        '<?xml version="1.0" encoding="utf-8"?>\n'
        '<X3D version="3.3" xmlns:xsd="http://www.w3.org/2001/XMLSchema-instance">\n'
        ' <Scene>\n  <Shape>\n'
        f'   <IndexedFaceSet coordIndex="{idx_str}" solid="false">\n'
        f'    <Coordinate point="{coord_str}"/>\n'
        '   </IndexedFaceSet>\n  </Shape>\n </Scene>\n</X3D>'
    )
    Path(dst).write_text(x3d, encoding='utf-8')

def _mesh_to_vrml(src, dst, **kw):
    """Mesh → VRML/WRL (Virtual Reality Modeling Language 2.0)"""
    import trimesh
    m = trimesh.load(src, force="mesh")
    verts, faces = m.vertices, m.faces
    pts = ",\n        ".join(f"{v[0]:.4f} {v[1]:.4f} {v[2]:.4f}" for v in verts)
    idx = ",\n        ".join(f"{f[0]} {f[1]} {f[2]} -1" for f in faces)
    vrml = (
        "#VRML V2.0 utf8\n"
        "Shape {\n  geometry IndexedFaceSet {\n"
        "    coord Coordinate {\n      point [\n"
        f"        {pts}\n      ]\n    }}\n"
        "    coordIndex [\n"
        f"        {idx}\n    ]\n  }}\n}}"
    )
    Path(dst).write_text(vrml, encoding='utf-8')


def _rhino_to_mesh(src, dst, **kw):
    """Rhino 3DM → STL/OBJ/PLY via rhino3dm"""
    import rhino3dm
    import trimesh, numpy as np
    model = rhino3dm.File3dm.Read(src)
    all_verts, all_faces = [], []
    offset = 0
    for obj in model.Objects:
        geo = obj.Geometry
        if isinstance(geo, rhino3dm.Mesh):
            verts = [(v.X, v.Y, v.Z) for v in geo.Vertices]
            faces = [(f[0]+offset, f[1]+offset, f[2]+offset) for f in geo.Faces]
            all_verts.extend(verts)
            all_faces.extend(faces)
            offset += len(verts)
        elif hasattr(geo, "ToMesh"):
            try:
                mesh = geo.ToMesh(rhino3dm.MeshingParameters.Default)
                if mesh:
                    verts = [(v.X, v.Y, v.Z) for v in mesh.Vertices]
                    faces = [(f[0]+offset, f[1]+offset, f[2]+offset) for f in mesh.Faces]
                    all_verts.extend(verts)
                    all_faces.extend(faces)
                    offset += len(verts)
            except Exception:
                pass
    if not all_verts:
        raise ValueError("Keine Mesh-Geometrie in der Rhino-Datei gefunden")
    m = trimesh.Trimesh(vertices=np.array(all_verts), faces=np.array(all_faces))
    m.export(dst)

def _blend_export(src, dst, export_fmt="GLB", **kw):
    """Blender headless Export — .blend/FBX/MA → GLB/STL/OBJ/etc."""
    if not _BLENDER:
        raise RuntimeError("Blender ist nicht installiert (apt install blender)")
    ext = norm(Path(dst).suffix)
    dst_abs = str(Path(dst).resolve())
    src_abs = str(Path(src).resolve())
    src_fmt = norm(Path(src).suffix)

    if src_fmt == "blend":
        open_cmd = f"bpy.ops.wm.open_mainfile(filepath=r'{src_abs}')"
    elif src_fmt == "fbx":
        open_cmd = f"bpy.ops.wm.read_factory_settings(use_empty=True); bpy.ops.import_scene.fbx(filepath=r'{src_abs}')"
    elif src_fmt in ("obj", "mtl"):
        open_cmd = f"bpy.ops.wm.read_factory_settings(use_empty=True); bpy.ops.wm.obj_import(filepath=r'{src_abs}')"
    elif src_fmt == "stl":
        open_cmd = f"bpy.ops.wm.read_factory_settings(use_empty=True); bpy.ops.wm.stl_import(filepath=r'{src_abs}')"
    elif src_fmt == "dae":
        open_cmd = f"bpy.ops.wm.read_factory_settings(use_empty=True); bpy.ops.wm.collada_import(filepath=r'{src_abs}')"
    elif src_fmt == "glb":
        open_cmd = f"bpy.ops.wm.read_factory_settings(use_empty=True); bpy.ops.import_scene.gltf(filepath=r'{src_abs}')"
    else:
        raise ValueError(f"Blender-Import für '{src_fmt}' nicht implementiert")

    if ext in ("glb", "gltf"):
        export_cmd = f"bpy.ops.export_scene.gltf(filepath=r'{dst_abs}', export_format='GLB')"
    elif ext == "fbx":
        export_cmd = f"bpy.ops.export_scene.fbx(filepath=r'{dst_abs}')"
    elif ext == "stl":
        export_cmd = f"bpy.ops.wm.stl_export(filepath=r'{dst_abs}')"
    elif ext == "obj":
        export_cmd = f"bpy.ops.wm.obj_export(filepath=r'{dst_abs}')"
    elif ext == "dae":
        export_cmd = f"bpy.ops.wm.collada_export(filepath=r'{dst_abs}')"
    elif ext == "usdz":
        export_cmd = f"bpy.ops.wm.usd_export(filepath=r'{dst_abs}')"
    elif ext == "abc":
        export_cmd = f"bpy.ops.wm.alembic_export(filepath=r'{dst_abs}')"
    elif ext == "ply":
        export_cmd = f"bpy.ops.export_mesh.ply(filepath=r'{dst_abs}')"
    elif ext == "x3d":
        export_cmd = f"bpy.ops.export_scene.x3d(filepath=r'{dst_abs}')"
    else:
        raise ValueError(f"Blender-Export für '{ext}' nicht implementiert")

    script = f"import bpy\n{open_cmd}\n{export_cmd}\n"
    with tempfile.NamedTemporaryFile(mode="w", suffix=".py", delete=False) as f:
        f.write(script)
        script_path = f.name
    try:
        r = _run(["blender", "--background", "--python", script_path], timeout=600)
        if r.returncode != 0 and not Path(dst).exists():
            raise RuntimeError(f"Blender: {r.stderr.decode(errors='replace')[-400:]}")
    finally:
        try: os.unlink(script_path)
        except: pass

# ══════════ CAD / BREP ══════════

def _step_to_stl(src, dst, **kw):
    """STEP → STL via OpenCascade (OCP)"""
    from OCP.STEPControl import STEPControl_Reader
    from OCP.BRepMesh import BRepMesh_IncrementalMesh
    from OCP.StlAPI import StlAPI_Writer
    from OCP.IFSelect import IFSelect_RetDone
    reader = STEPControl_Reader()
    if reader.ReadFile(src) != IFSelect_RetDone:
        raise RuntimeError("STEP-Datei konnte nicht gelesen werden")
    reader.TransferRoots()
    shape = reader.OneShape()
    BRepMesh_IncrementalMesh(shape, 0.01, False, 0.1).Perform()
    StlAPI_Writer().Write(shape, dst)

def _step_to_brep(src, dst, **kw):
    """STEP → BREP via OpenCascade (OCP)"""
    from OCP.STEPControl import STEPControl_Reader
    from OCP.BRepTools import BRepTools
    from OCP.IFSelect import IFSelect_RetDone
    reader = STEPControl_Reader()
    if reader.ReadFile(src) != IFSelect_RetDone:
        raise RuntimeError("STEP-Datei konnte nicht gelesen werden")
    reader.TransferRoots()
    shape = reader.OneShape()
    BRepTools.Write_s(shape, dst)

def _step_to_iges(src, dst, **kw):
    from OCP.STEPControl import STEPControl_Reader
    from OCP.IGESControl import IGESControl_Writer
    r = STEPControl_Reader(); r.ReadFile(src); r.TransferRoots()
    w = IGESControl_Writer(); w.AddShape(r.Shape()); w.Write(dst)

def _brep_to_step(src, dst, **kw):
    """BREP → STEP via OpenCascade (OCP)"""
    from OCP.BRep import BRep_Builder
    from OCP.BRepTools import BRepTools
    from OCP.TopoDS import TopoDS_Shape
    from OCP.STEPControl import STEPControl_Writer, STEPControl_AsIs
    builder = BRep_Builder()
    shape = TopoDS_Shape()
    BRepTools.Read_s(shape, src, builder)
    w = STEPControl_Writer()
    w.Transfer(shape, STEPControl_AsIs)
    w.Write(dst)

def _brep_to_stl(src, dst, **kw):
    """BREP → STL via OpenCascade (OCP)"""
    from OCP.BRep import BRep_Builder
    from OCP.BRepTools import BRepTools
    from OCP.TopoDS import TopoDS_Shape
    from OCP.BRepMesh import BRepMesh_IncrementalMesh
    from OCP.StlAPI import StlAPI_Writer
    builder = BRep_Builder()
    shape = TopoDS_Shape()
    BRepTools.Read_s(shape, src, builder)
    BRepMesh_IncrementalMesh(shape, 0.01, False, 0.1).Perform()
    StlAPI_Writer().Write(shape, dst)

def _iges_to_step(src, dst, **kw):
    from OCP.IGESControl import IGESControl_Reader
    from OCP.STEPControl import STEPControl_Writer, STEPControl_AsIs
    r = IGESControl_Reader(); r.ReadFile(src); r.TransferRoots()
    w = STEPControl_Writer(); w.Transfer(r.Shape(), STEPControl_AsIs); w.Write(dst)

def _stl_to_step(src, dst, **kw):
    from OCP.StlAPI import StlAPI_Reader
    from OCP.TopoDS import TopoDS_Shape
    from OCP.STEPControl import STEPControl_Writer, STEPControl_AsIs
    from OCP.BRepBuilderAPI import BRepBuilderAPI_Sewing
    from OCP.ShapeFix import ShapeFix_Shape
    reader = StlAPI_Reader(); shape = TopoDS_Shape(); reader.Read(shape, src)
    sew = BRepBuilderAPI_Sewing(0.01); sew.Add(shape); sew.Perform(); sewn = sew.SewedShape()
    try: fixer = ShapeFix_Shape(sewn); fixer.Perform(); sewn = fixer.Shape()
    except Exception: pass
    writer = STEPControl_Writer(); writer.Transfer(sewn, STEPControl_AsIs); writer.Write(dst)


def _dxf_to_svg(src, dst, **kw):
    import ezdxf
    from ezdxf.addons.drawing import RenderContext, Frontend
    from ezdxf.addons.drawing.svg import SVGBackend
    doc = ezdxf.readfile(src)
    ctx = RenderContext(doc)
    backend = SVGBackend()
    Frontend(ctx, backend).draw_layout(doc.modelspace())
    try:
        svg_str = backend.get_string()
    except TypeError:
        from ezdxf.addons.drawing.layout import Page, Units
        svg_str = backend.get_string(Page(297, 210, Units.mm))
    Path(dst).write_text(svg_str, encoding="utf-8")

def _dxf_to_png(src, dst, **kw):
    """DXF → PNG via ezdxf + matplotlib"""
    import ezdxf
    from ezdxf.addons.drawing import RenderContext, Frontend
    from ezdxf.addons.drawing.matplotlib import MatplotlibBackend
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    doc = ezdxf.readfile(src)
    fig = plt.figure(figsize=(12, 9))
    ax = fig.add_axes([0, 0, 1, 1])
    ctx = RenderContext(doc)
    out_be = MatplotlibBackend(ax)
    Frontend(ctx, out_be).draw_layout(doc.modelspace())
    ax.set_axis_off()
    fig.savefig(dst, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)

def _dxf_to_pdf(src, dst, **kw):
    """DXF → PDF via ezdxf + matplotlib"""
    import ezdxf
    from ezdxf.addons.drawing import RenderContext, Frontend
    from ezdxf.addons.drawing.matplotlib import MatplotlibBackend
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    doc = ezdxf.readfile(src)
    fig = plt.figure(figsize=(11.69, 8.27))
    ax = fig.add_axes([0, 0, 1, 1])
    ctx = RenderContext(doc)
    out_be = MatplotlibBackend(ax)
    Frontend(ctx, out_be).draw_layout(doc.modelspace())
    ax.set_axis_off()
    fig.savefig(dst, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)

def _dxf_to_txt(src, dst, **kw):
    """DXF → TXT (Entitätsliste mit Typ-Statistik)"""
    import ezdxf
    doc = ezdxf.readfile(src)
    msp = doc.modelspace()
    counts: dict = {}
    for entity in msp:
        t = entity.dxftype()
        counts[t] = counts.get(t, 0) + 1
    lines = [f"DXF-Datei: {Path(src).name}"]
    try:
        lines.append(f"Einheiten: {doc.header.get('$INSUNITS', 'unbekannt')}")
    except Exception:
        pass
    lines += ["", "Entitaeten:"] + [f"  {t}: {n}" for t, n in sorted(counts.items())]
    Path(dst).write_text("\n".join(lines), encoding="utf-8")

def _dxf_to_mesh(src, dst, **kw):
    """DXF → STL/OBJ/PLY/GLB: 3D-Entities (3DFACE/MESH) oder 2D-Konturen als flaches Mesh"""
    import ezdxf
    import trimesh
    import numpy as np
    doc = ezdxf.readfile(src)
    verts, faces = [], []
    vi = 0

    # ── Schritt 1: echte 3D-Entities ──────────────────────────────────────────
    for e in doc.modelspace():
        if e.dxftype() == "3DFACE":
            p = [e.dxf.vtx0, e.dxf.vtx1, e.dxf.vtx2, e.dxf.vtx3]
            v = [(pt.x, pt.y, pt.z) for pt in p]
            if (v[2][0] == v[3][0] and v[2][1] == v[3][1] and v[2][2] == v[3][2]):
                verts.extend(v[:3]); faces.append([vi, vi+1, vi+2]); vi += 3
            else:
                verts.extend(v)
                faces.append([vi, vi+1, vi+2]); faces.append([vi, vi+2, vi+3]); vi += 4
        elif e.dxftype() in ("MESH", "POLYFACE"):
            try:
                mb = e.mesh_builder() if hasattr(e, "mesh_builder") else None
                if mb:
                    base = vi
                    verts.extend(mb.vertices)
                    faces.extend([[f[0]+base, f[1]+base, f[2]+base] for f in mb.faces])
                    vi += len(mb.vertices)
            except Exception:
                pass

    # ── Schritt 2: 2D-Fallback — alle 2D-Entities als flaches Mesh (z=0) ─────
    if not faces:
        import math as _math

        def _entity_to_pts(e):
            """Konvertiert ein DXF-Entity zu einer Liste von (x, y) Punkten."""
            t = e.dxftype()
            try:
                if t == "LWPOLYLINE":
                    return [(p[0], p[1]) for p in e.get_points("xy")]
                elif t == "CIRCLE":
                    cx, cy = e.dxf.center.x, e.dxf.center.y
                    r = e.dxf.radius
                    n = max(16, min(128, int(2 * _math.pi * r)))
                    return [(cx + r * _math.cos(2*_math.pi*i/n),
                             cy + r * _math.sin(2*_math.pi*i/n)) for i in range(n)]
                elif t == "ARC":
                    cx, cy = e.dxf.center.x, e.dxf.center.y
                    r = e.dxf.radius
                    a0 = _math.radians(e.dxf.start_angle)
                    a1 = _math.radians(e.dxf.end_angle)
                    if a1 < a0: a1 += 2 * _math.pi
                    span = a1 - a0
                    n = max(4, int(span * r))
                    # Bogen + Mittelpunkt = geschlossenes Dreiecksfächer-Polygon
                    arc_pts = [(cx + r * _math.cos(a0 + span*i/n),
                                cy + r * _math.sin(a0 + span*i/n)) for i in range(n+1)]
                    return [(cx, cy)] + arc_pts
                elif t == "SPLINE":
                    pts = list(e.flattening(0.1))
                    return [(p[0], p[1]) for p in pts]
                elif t == "POLYLINE":
                    return [(p[0], p[1]) for p in e.points()]
                elif t in ("ELLIPSE",):
                    pts = list(e.flattening(0.1))
                    return [(p[0], p[1]) for p in pts]
                elif t == "HATCH":
                    # HATCH-Boundaries sind perfekte geschlossene Konturen
                    result = []
                    for bp in e.paths:
                        pts = []
                        if hasattr(bp, "vertices"):          # PolylinePath
                            pts = [(v[0], v[1]) for v in bp.vertices]
                        elif hasattr(bp, "edges"):           # EdgePath
                            for edge in bp.edges:
                                et = type(edge).__name__
                                if et == "LineEdge":
                                    pts.append((edge.start[0], edge.start[1]))
                                elif et == "ArcEdge":
                                    ecx, ecy = edge.center
                                    er = edge.radius
                                    ea0 = _math.radians(edge.start_angle)
                                    ea1 = _math.radians(edge.end_angle)
                                    if ea1 < ea0: ea1 += 2*_math.pi
                                    for i in range(8):
                                        a = ea0 + (ea1-ea0)*i/8
                                        pts.append((ecx + er*_math.cos(a), ecy + er*_math.sin(a)))
                        if len(pts) >= 3:
                            result.append(pts)
                    return result[0] if result else []
            except Exception:
                pass
            return []

        def _collect_paths(msp, depth=0):
            """Sammelt 2D-Pfade aus Entities, expandiert INSERT-Blöcke rekursiv."""
            result = []
            if depth > 3:
                return result
            for e in msp:
                pts = _entity_to_pts(e)
                if len(pts) >= 3:
                    result.append(pts)
                elif e.dxftype() == "INSERT":
                    try:
                        block = doc.blocks[e.dxf.name]
                        result.extend(_collect_paths(block, depth + 1))
                    except Exception:
                        pass
            return result

        all_paths = _collect_paths(doc.modelspace())

        # Fächer-Triangulierung: funktioniert für konvexe und schwach konkave Polygone
        for pts2d in all_paths:
            if len(pts2d) < 3:
                continue
            p0 = (pts2d[0][0], pts2d[0][1], 0.0)
            for i in range(1, len(pts2d) - 1):
                p1 = (pts2d[i][0],   pts2d[i][1],   0.0)
                p2 = (pts2d[i+1][0], pts2d[i+1][1], 0.0)
                if p0 != p1 and p1 != p2 and p0 != p2:
                    verts.extend([p0, p1, p2])
                    faces.append([vi, vi+1, vi+2]); vi += 3

    if not faces:
        raise RuntimeError("Keine Geometrie zum Konvertieren gefunden.")
    mesh = trimesh.Trimesh(vertices=np.array(verts, dtype=float),
                           faces=np.array(faces, dtype=int))
    mesh.export(dst)

def _mesh_to_dxf(src, dst, **kw):
    """STL/OBJ/PLY/GLB → DXF: Mesh-Flächen als 3DFACE-Entitäten via ezdxf"""
    import trimesh
    import ezdxf
    mesh = trimesh.load(src, force="mesh")
    doc = ezdxf.new("R2010")
    msp = doc.modelspace()
    for tri in mesh.faces:
        v = mesh.vertices[tri]
        p = [tuple(float(x) for x in v[i]) for i in range(3)]
        msp.add_3dface([p[0], p[1], p[2], p[2]])
    doc.saveas(dst)

def _svg_to_dxf(src, dst, **kw):
    """SVG → DXF via Inkscape (--export-type=dxf)"""
    if not _INKSCAPE:
        raise RuntimeError("SVG→DXF erfordert Inkscape (apt install inkscape)")
    r = _run(["inkscape", "--export-type=dxf", f"--export-filename={dst}", src])
    if r.returncode != 0 or not Path(dst).exists():
        raise RuntimeError(f"Inkscape SVG→DXF: {r.stderr.decode(errors='replace')[-300:]}")

def _ifc_to_obj(src, dst, **kw):
    """IFC → OBJ via ifcopenshell.geom"""
    import ifcopenshell
    import ifcopenshell.geom
    import trimesh
    import numpy as np
    model = ifcopenshell.open(src)
    settings = ifcopenshell.geom.settings()
    settings.set(settings.USE_WORLD_COORDS, True)
    vertices_all, faces_all = [], []
    offset = 0
    for product in model.by_type("IfcProduct"):
        try:
            shape = ifcopenshell.geom.create_shape(settings, product)
            geom = shape.geometry
            verts = np.array(geom.verts).reshape(-1, 3)
            faces = np.array(geom.faces).reshape(-1, 3)
            vertices_all.append(verts)
            faces_all.append(faces + offset)
            offset += len(verts)
        except Exception:
            pass
    if not vertices_all:
        raise ValueError("Keine Geometrie im IFC gefunden")
    mesh = trimesh.Trimesh(
        vertices=np.vstack(vertices_all),
        faces=np.vstack(faces_all)
    )
    mesh.export(dst)

def _ifc_to_glb(src, dst, **kw):
    # Versuch: IfcConvert CLI
    try:
        r = _run(["IfcConvert", src, dst], timeout=300)
        if r.returncode == 0 and Path(dst).exists(): return
    except FileNotFoundError: pass
    # Fallback via OBJ
    obj_tmp = dst + "_tmp.obj"
    try:
        _ifc_to_obj(src, obj_tmp)
        import trimesh
        trimesh.load(obj_tmp).export(dst)
    finally:
        try: os.unlink(obj_tmp)
        except: pass

# ══════════ GIS / GEODATEN ══════════

def _shp_to_geojson(src, dst, **kw):
    import shapefile, json
    sf = shapefile.Reader(src)
    features = []
    for shape_rec in sf.shapeRecords():
        features.append({
            "type": "Feature",
            "geometry": shape_rec.shape.__geo_interface__,
            "properties": dict(zip([f[0] for f in sf.fields[1:]], shape_rec.record)),
        })
    Path(dst).write_text(json.dumps({
        "type": "FeatureCollection", "features": features
    }, indent=2, default=str), encoding="utf-8")

def _geojson_to_kml(src, dst, **kw):
    import json
    data = json.loads(Path(src).read_text(encoding="utf-8"))
    lines = ['<?xml version="1.0" encoding="UTF-8"?>',
             '<kml xmlns="http://www.opengis.net/kml/2.2"><Document>']
    for feat in data.get("features", [data]):
        geo = feat.get("geometry", feat)
        props = feat.get("properties", {})
        name = props.get("name", "Feature") if props else "Feature"
        lines.append(f"<Placemark><name>{name}</name>")
        gtype = geo.get("type","")
        coords = geo.get("coordinates",[])
        if gtype == "Point":
            c = coords
            lines.append(f"<Point><coordinates>{c[0]},{c[1]}</coordinates></Point>")
        elif gtype in ("LineString", "MultiLineString"):
            if gtype == "LineString": coords = [coords]
            for ring in coords:
                coord_str = " ".join(f"{c[0]},{c[1]}" for c in ring)
                lines.append(f"<LineString><coordinates>{coord_str}</coordinates></LineString>")
        elif gtype in ("Polygon", "MultiPolygon"):
            if gtype == "Polygon": coords = [coords]
            for poly in coords:
                ring = poly[0] if poly else []
                coord_str = " ".join(f"{c[0]},{c[1]}" for c in ring)
                lines.append(f"<Polygon><outerBoundaryIs><LinearRing><coordinates>{coord_str}</coordinates></LinearRing></outerBoundaryIs></Polygon>")
        lines.append("</Placemark>")
    lines.append("</Document></kml>")
    Path(dst).write_text("\n".join(lines), encoding="utf-8")

def _kml_to_geojson(src, dst, **kw):
    import re, json
    kml = Path(src).read_text(encoding="utf-8", errors="replace")
    features = []
    for coord_match in re.finditer(r'<coordinates>(.*?)</coordinates>', kml, re.DOTALL):
        raw = coord_match.group(1).strip()
        pairs = []
        for token in raw.split():
            parts = token.split(",")
            if len(parts) >= 2:
                try: pairs.append([float(parts[0]), float(parts[1])])
                except ValueError: pass
        if pairs:
            geo = {"type": "Point", "coordinates": pairs[0]} if len(pairs)==1 else \
                  {"type": "LineString", "coordinates": pairs}
            features.append({"type": "Feature", "geometry": geo, "properties": {}})
    Path(dst).write_text(json.dumps({"type":"FeatureCollection","features":features},indent=2),
                         encoding="utf-8")

def _gpx_to_geojson(src, dst, **kw):
    import gpxpy, json
    gpx = gpxpy.parse(Path(src).read_text(encoding="utf-8"))
    features = []
    for track in gpx.tracks:
        for seg in track.segments:
            coords = [[p.longitude, p.latitude] for p in seg.points]
            features.append({"type":"Feature","geometry":{"type":"LineString","coordinates":coords},
                             "properties":{"name":track.name}})
    for wpt in gpx.waypoints:
        features.append({"type":"Feature","geometry":{"type":"Point","coordinates":[wpt.longitude,wpt.latitude]},
                        "properties":{"name":wpt.name}})
    Path(dst).write_text(json.dumps({"type":"FeatureCollection","features":features},indent=2),
                         encoding="utf-8")

def _gpx_to_csv(src, dst, **kw):
    import gpxpy, csv
    gpx = gpxpy.parse(Path(src).read_text(encoding="utf-8"))
    rows = []
    for track in gpx.tracks:
        for seg in track.segments:
            for p in seg.points:
                rows.append({"lat":p.latitude,"lon":p.longitude,"ele":p.elevation,"time":str(p.time),"track":track.name})
    with open(dst,"w",newline="",encoding="utf-8") as f:
        if rows:
            writer = csv.DictWriter(f, fieldnames=list(rows[0].keys()))
            writer.writeheader(); writer.writerows(rows)

def _geojson_to_shp(src, dst, **kw):
    import shapefile, json
    data = json.loads(Path(src).read_text(encoding="utf-8"))
    features = data.get("features", [])
    w = shapefile.Writer(dst)
    w.autoBalance = 1
    try:
        if features:
            props = features[0].get("properties",{}) or {}
            if props:
                for k, v in props.items():
                    w.field(str(k)[:10], "C", 100)
            else:
                w.field("ID", "N", 10)
            for i, feat in enumerate(features):
                geo = feat.get("geometry", {})
                gtype = geo.get("type", "")
                fprops = feat.get("properties", {}) or {}
                vals = [str(v) for v in fprops.values()] if fprops else [i]
                if gtype == "Point":
                    w.point(*geo["coordinates"][:2]); w.record(*vals)
                elif gtype == "LineString":
                    w.line([geo["coordinates"]]); w.record(*vals)
                elif gtype == "Polygon":
                    w.poly(geo["coordinates"]); w.record(*vals)
    finally:
        try:
            w.close()
        except Exception:
            pass  # suppress Windows file-lock errors from shapefile __del__

def _geojson_to_shp_fiona(src, dst, **kw):
    import fiona, json
    from fiona.crs import from_epsg
    data = json.loads(Path(src).read_text(encoding="utf-8"))
    features = data.get("features", [])
    if not features:
        raise ValueError("GeoJSON enthält keine Features")
    gtype = features[0].get("geometry", {}).get("type", "Point")
    # Fiona erwartet ein Verzeichnis als Ziel für Shapefiles
    dst_dir = str(Path(dst).parent)
    dst_name = Path(dst).stem
    props = features[0].get("properties") or {}
    schema = {
        "geometry": gtype,
        "properties": {k: "str" for k in props.keys()} if props else {"ID": "int"},
    }
    with fiona.open(Path(dst_dir) / (dst_name + ".shp"), "w",
                    driver="ESRI Shapefile", schema=schema,
                    crs=from_epsg(4326)) as out:
        for i, feat in enumerate(features):
            geo = feat.get("geometry")
            fprops = feat.get("properties") or {}
            if not geo: continue
            rec_props = {k: str(v) for k, v in fprops.items()} if props else {"ID": i}
            out.write({"geometry": geo, "properties": rec_props})

# ══════════ FONTS ══════════

def _font_convert(src, dst, **kw):
    """TTF ↔ OTF ↔ WOFF ↔ WOFF2 via fonttools"""
    from fontTools.ttLib import TTFont
    font = TTFont(src)
    ext = norm(Path(dst).suffix)
    if ext == "woff":
        font.flavor = "woff"
    elif ext == "woff2":
        font.flavor = "woff2"
    else:
        font.flavor = None
    font.save(dst)

# ══════════ STICKEREI ══════════

def _embroidery_convert(src, dst, **kw):
    """Alle Stickerei-Formate via pyembroidery (80+ Formate)"""
    import pyembroidery
    pattern = pyembroidery.read(src)
    if pattern is None:
        raise ValueError(f"Stickerei-Datei konnte nicht gelesen werden: {src}")
    pyembroidery.write(pattern, dst)

def _embroidery_to_csv(src, dst, **kw):
    import pyembroidery, csv
    pattern = pyembroidery.read(src)
    if pattern is None: raise ValueError("Stickerei-Datei unlesbar")
    with open(dst,"w",newline="",encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(["x","y","command"])
        for stitch in pattern.stitches:
            writer.writerow([stitch[0], stitch[1], stitch[2]])

def _embroidery_to_svg(src, dst, **kw):
    import pyembroidery
    pattern = pyembroidery.read(src)
    if pattern is None: raise ValueError("Stickerei-Datei unlesbar")
    pyembroidery.write(pattern, dst)

# ══════════ MEDIZIN ══════════

def _dcm_to_png(src, dst, **kw):
    """DICOM → PNG/JPEG via pydicom"""
    import pydicom
    from PIL import Image
    import numpy as np
    ds = pydicom.dcmread(src)
    arr = ds.pixel_array.astype(float)
    arr = ((arr - arr.min()) / (arr.max() - arr.min() + 1e-9) * 255).astype(np.uint8)
    img = Image.fromarray(arr)
    if norm(Path(dst).suffix) == "jpeg":
        img = img.convert("RGB")
    img.save(dst)

def _dcm_to_nii(src, dst, **kw):
    """DICOM → NIfTI via nibabel/SimpleITK"""
    try:
        import SimpleITK as sitk
        image = sitk.ReadImage(src)
        sitk.WriteImage(image, dst)
        return
    except ImportError: pass
    raise RuntimeError("pip install SimpleITK  für DICOM→NIfTI")

# ══════════ LIDAR / PUNKTWOLKE ══════════

def _las_to_csv(src, dst, **kw):
    import laspy, csv
    las = laspy.read(src)
    with open(dst,"w",newline="",encoding="utf-8") as f:
        writer = csv.writer(f)
        writer.writerow(["x","y","z"])
        for x,y,z in zip(las.x, las.y, las.z):
            writer.writerow([float(x),float(y),float(z)])

def _las_to_ply(src, dst, **kw):
    import laspy, trimesh
    import numpy as np
    las = laspy.read(src)
    pts = np.column_stack([las.x, las.y, las.z])
    cloud = trimesh.points.PointCloud(pts)
    cloud.export(dst)

def _laz_to_las(src, dst, **kw):
    import laspy
    las = laspy.read(src)
    las.write(dst)

# ══════════ PCB / GERBER ══════════

def _gerber_to_svg(src, dst, **kw):
    """Gerber → SVG via gerber-tools"""
    import gerber
    from gerber.render.svg_backend import GerberSvgBackend
    layer = gerber.load(src)
    settings = gerber.render.RenderSettings()
    backend = GerberSvgBackend(margin=2.0)
    ctx = gerber.render.GerberCairoContext(scale=300)
    ctx.render_layer(layer, settings=settings, bgsettings=None)
    ctx.dump(dst)

def _gerber_to_png(src, dst, **kw):
    """Gerber → PNG via gerber-tools"""
    import gerber
    from gerber.render import GerberCairoContext, RenderSettings
    layer = gerber.load(src)
    ctx = GerberCairoContext(scale=300)
    ctx.render_layer(layer, settings=RenderSettings())
    ctx.dump(dst)

# ══════════ G-CODE ══════════

def _gcode_to_csv(src, dst, **kw):
    """G-Code → CSV (Maschinenkoordinaten)"""
    import re, csv
    lines = Path(src).read_text(encoding="utf-8", errors="replace").split("\n")
    rows = []
    x = y = z = 0.0
    for line in lines:
        code = line.split(";")[0].strip().upper()
        if code.startswith("G0") or code.startswith("G1"):
            for match in re.finditer(r"([XYZ])([-\d.]+)", code):
                ax, val = match.group(1), float(match.group(2))
                if ax == "X": x = val
                elif ax == "Y": y = val
                elif ax == "Z": z = val
            rows.append({"cmd": code[:3].strip(), "x":x,"y":y,"z":z})
    with open(dst,"w",newline="",encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=["cmd","x","y","z"])
        writer.writeheader(); writer.writerows(rows)

def _gcode_to_svg(src, dst, **kw):
    """G-Code → SVG (2D-Toolpath-Visualisierung)"""
    import re
    lines = Path(src).read_text(encoding="utf-8", errors="replace").split("\n")
    paths, x, y = [], 0.0, 0.0
    cur_path = []
    for line in lines:
        code = line.split(";")[0].strip().upper()
        if code.startswith("G0") or code.startswith("G1"):
            for match in re.finditer(r"([XY])([-\d.]+)", code):
                ax,val = match.group(1), float(match.group(2))
                if ax=="X": x=val
                else: y=val
            mode = "G0" if code.startswith("G0") else "G1"
            if mode == "G0" and cur_path:
                paths.append(("move", cur_path)); cur_path = []
            cur_path.append((x, y))
    if cur_path: paths.append(("cut", cur_path))
    all_pts = [p for _, pts in paths for p in pts]
    if not all_pts:
        Path(dst).write_text('<svg xmlns="http://www.w3.org/2000/svg"/>', encoding="utf-8"); return
    xs, ys = [p[0] for p in all_pts], [p[1] for p in all_pts]
    minx,miny,maxx,maxy = min(xs),min(ys),max(xs),max(ys)
    W,H = max(maxx-minx,1), max(maxy-miny,1)
    lines_svg = [f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="{minx} {miny} {W} {H}" width="800" height="600">',
                 f'<rect width="100%" height="100%" fill="#111"/>']
    for ptype, pts in paths:
        color = "#0f0" if ptype=="cut" else "#555"
        d = "M " + " L ".join(f"{p[0]} {p[1]}" for p in pts)
        lines_svg.append(f'<path d="{d}" stroke="{color}" stroke-width="{W*0.001}" fill="none"/>')
    lines_svg.append("</svg>")
    Path(dst).write_text("\n".join(lines_svg), encoding="utf-8")

# ══════════ ARCHIVE ══════════

def _to_zip(src, dst, **kw):
    import zipfile
    with zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as z:
        z.write(src, Path(src).name)

def _zip_to_tar(src, dst, **kw):
    import zipfile, tarfile
    with tempfile.TemporaryDirectory() as td:
        with zipfile.ZipFile(src) as z: z.extractall(td)
        mode = "w:gz" if dst.endswith(".gz") else "w:bz2" if dst.endswith(".bz2") else "w:xz" if dst.endswith(".xz") else "w"
        with tarfile.open(dst, mode) as t: t.add(td, arcname="")

def _to_7z(src, dst, **kw):
    import py7zr
    with py7zr.SevenZipFile(dst, "w") as z:
        z.write(src, Path(src).name)

def _7z_extract_first(src, dst, **kw):
    import py7zr
    with py7zr.SevenZipFile(src, "r") as z:
        names = z.getnames()
        if not names: raise ValueError("7z ist leer")
        with tempfile.TemporaryDirectory() as td:
            z.extractall(td)
            shutil.copy(os.path.join(td, names[0]), dst)

def _rar_extract_first(src, dst, **kw):
    import rarfile
    with rarfile.RarFile(src) as rf:
        names = rf.namelist()
        if not names: raise ValueError("RAR ist leer")
        rf.extract(names[0], Path(dst).parent)
        extracted = Path(dst).parent / names[0]
        if str(extracted) != dst: shutil.move(str(extracted), dst)

# ── Konvertierungsgraph ───────────────────────────────────────────────────────

GRAPH: dict = {}

_IMG_FMTS   = ["jpeg","png","gif","bmp","tiff","webp","ico","tga","pcx","ppm","pgm","pbm"]
_RAW_FMTS   = ["cr3","cr2","nef","nrw","arw","srf","sr2","raf","orf","rw2","dng","pef","3fr","fff","mrw","erf","mos","mef"]
_AUDIO_FMTS = ["mp3","wav","ogg","flac","aac","opus","wma","aiff","ac3","mp2"]
_VIDEO_FMTS = ["mp4","avi","mov","mkv","webm","flv","wmv","mpg","ts","3gp","mxf","rm"]
_DATA_FMTS  = ["csv","tsv","json","xlsx","xls","xml","parquet","feather","orc","toml"]
_MESH_FMTS  = ["stl","obj","ply","3mf","glb","off","dae","amf","vrml"]
_EMBD_FMTS  = ["pes","dst","exp","jef","xxx","vp3","hus","sew","vip","csd","esd","u01","phb","phc","inb"]
_FONT_FMTS  = ["ttf","otf","woff","woff2"]
_GIS_FMTS   = ["geojson","kml","kmz","gpx","shp"]
_GERBER_FMTS= ["gbr","gtl","gbl","gts","gbs","gto","gbo","gko","ger","art"]
_CAD_FMTS   = ["step","brep","iges"]

# ── Bilder ────────────────────────────────────────────────────────────────────
if _PIL:
    for a in _IMG_FMTS:
        for b in _IMG_FMTS:
            if a != b: GRAPH[(a,b)] = _img_to_img
    for a in _IMG_FMTS:
        GRAPH[(a,"pdf")] = _img_to_pdf
    GRAPH[("pdf","jpeg")] = _pdf_to_img
    GRAPH[("pdf","png")] = _pdf_to_img

if _PIL and _PILHEIF:
    for b in _IMG_FMTS + ["pdf"]:
        GRAPH[("heic",b)] = _heic_to_img
        GRAPH[("heif",b)] = _heic_to_img
        GRAPH[("avif",b)] = _heic_to_img

if _RAWPY:
    for a in _RAW_FMTS:
        for b in ["jpeg","png","tiff","bmp","webp"]:
            GRAPH[(a,b)] = _raw_to_img

if _GHOSTSCR:
    GRAPH[("eps","png")] = _eps_to_png
    GRAPH[("eps","jpeg")] = _eps_to_png

GRAPH[("svg","png")] = _svg_to_png

try:
    import psd_tools
    GRAPH[("psd","png")] = _psd_to_png
    GRAPH[("psb","png")] = _psd_to_png
except ImportError: pass

try:
    import pytesseract
    import platform as _plat
    _tess_ok = False
    if _plat.system() == "Windows":
        for _tp in [r"C:\Program Files\Tesseract-OCR\tesseract.exe",
                    r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe"]:
            if Path(_tp).exists():
                _tess_ok = True; break
    else:
        _tess_ok = True
    if _tess_ok:
        for a in ["jpeg","png","bmp","tiff","webp","gif","heic","heif","avif"]:
            GRAPH[(a,"txt")] = _img_to_txt_ocr
except ImportError:
    pass

# ── Dokumente ─────────────────────────────────────────────────────────────────
if _PYPDF or _PYMUPDF:
    GRAPH[("pdf","txt")] = _pdf_to_txt

if _REPORTLAB:
    GRAPH[("txt","pdf")] = _txt_to_pdf

if _DOCX:
    GRAPH[("docx","txt")]  = _docx_to_txt
    GRAPH[("docx","html")] = _docx_to_html
    GRAPH[("txt","docx")]  = _txt_to_docx

if _PDF2DOCX:
    GRAPH[("pdf","docx")] = _pdf_to_docx

if _PPTX:
    GRAPH[("pptx","pdf")]  = _pptx_to_pdf
    GRAPH[("pptx","html")] = _pptx_to_html
    GRAPH[("pptx","txt")]  = _pptx_to_txt

GRAPH[("markdown","html")] = _md_to_html
GRAPH[("html","txt")] = _html_to_txt
GRAPH[("txt","html")] = _txt_to_html

if _SOFFICE:
    for a in ["docx","doc","odt","rtf","pptx","ppt","odp","xlsx","xls","ods"]:
        GRAPH[(a,"pdf")] = _libreoffice
    GRAPH[("docx","odt")]  = _libreoffice
    GRAPH[("odt","docx")]  = _libreoffice
    GRAPH[("pptx","odp")]  = _libreoffice
    GRAPH[("odp","pptx")]  = _libreoffice
    GRAPH[("xlsx","ods")]  = _libreoffice
    GRAPH[("ods","xlsx")]  = _libreoffice

if _PANDOC:
    for a in ["markdown","rst","html","docx","epub","latex"]:
        for b in ["markdown","html","pdf","docx","epub","latex","txt"]:
            if a != b: GRAPH[(a,b)] = _pandoc
    GRAPH[("rst","html")] = _pandoc
    GRAPH[("latex","pdf")] = _pandoc

# Native fallback: pandoc→pdf needs LaTeX which is not installed on the server.
# Override docx→pdf with the python-docx + reportlab converter when LibreOffice
# is also unavailable (preserves text and headings without any system deps).
if not _SOFFICE:
    GRAPH[("docx","pdf")] = _docx_to_pdf_native

if _EBOOKLIB:
    GRAPH[("epub","txt")]  = _epub_to_txt
    GRAPH[("epub","html")] = _epub_to_html

if _NBFORMAT and _NBCONVERT:
    GRAPH[("ipynb","html")]     = _ipynb_to_html
    GRAPH[("ipynb","py")]       = _ipynb_to_py
    GRAPH[("ipynb","markdown")] = _ipynb_to_md

# ── Tabellen / Daten ──────────────────────────────────────────────────────────
if _PANDAS:
    for a in _DATA_FMTS:
        for b in _DATA_FMTS + ["html","markdown","txt"]:
            if a != b: GRAPH[(a,b)] = _data_convert
    GRAPH[("sqlite","json")] = _sqlite_to_json
    GRAPH[("db","json")]     = _sqlite_to_json
    GRAPH[("sqlite","csv")]  = _sqlite_to_csv
    GRAPH[("db","csv")]      = _sqlite_to_csv

# ── Audio ─────────────────────────────────────────────────────────────────────
if _FFMPEG:
    for a in _AUDIO_FMTS:
        for b in _AUDIO_FMTS:
            if a != b: GRAPH[(a,b)] = _audio_convert
    if _MIDO:
        GRAPH[("midi","json")] = _midi_to_json
        GRAPH[("mid","json")]  = _midi_to_json
        try:
            _run(["fluidsynth", "--version"], timeout=3)
            GRAPH[("midi","wav")] = _midi_to_wav
            GRAPH[("mid","wav")]  = _midi_to_wav
        except Exception: pass

# ── Video ─────────────────────────────────────────────────────────────────────
if _FFMPEG:
    for a in _VIDEO_FMTS:
        for b in _VIDEO_FMTS:
            if a != b: GRAPH[(a,b)] = _video_convert
        GRAPH[(a,"gif")] = _video_convert
        GRAPH[(a,"jpeg")] = _video_thumbnail
        GRAPH[(a,"png")] = _video_thumbnail
        for b in _AUDIO_FMTS:
            GRAPH[(a,b)] = _video_to_audio
    GRAPH[("gif","mp4")]  = _video_convert
    GRAPH[("gif","webm")] = _video_convert

# ── 3D Mesh ───────────────────────────────────────────────────────────────────
if _TRIMESH:
    for a in _MESH_FMTS:
        for b in _MESH_FMTS:
            if a != b: GRAPH[(a,b)] = _mesh_convert
    # AMF: eigener Writer (XML-basiert)
    for _a in _MESH_FMTS:
        if _a != "amf": GRAPH[(_a,"amf")] = _mesh_to_amf
    # 3MF: eigener Writer (ZIP-Paket, Fusion 360 kompatibel)
    for _a in _MESH_FMTS:
        if _a != "3mf": GRAPH[(_a,"3mf")] = _mesh_to_3mf
    # X3D: nur Export (trimesh liest X3D nicht nativ)
    for _a in _MESH_FMTS:
        GRAPH[(_a,"x3d")] = _mesh_to_x3d
    # VRML: eigener Writer; trimesh liest .wrl (via vrml-Alias)
    for _a in _MESH_FMTS:
        if _a != "vrml": GRAPH[(_a,"vrml")] = _mesh_to_vrml

if _RHINO3DM:
    for b in ["stl","obj","ply","glb","3mf"]:
        GRAPH[("3dm",b)] = _rhino_to_mesh

# ── CAD / BREP ────────────────────────────────────────────────────────────────
try:
    from OCP.STEPControl import STEPControl_Reader, STEPControl_Writer
    GRAPH[("step","stl")]  = _step_to_stl
    GRAPH[("step","brep")] = _step_to_brep
    GRAPH[("step","iges")] = _step_to_iges
    GRAPH[("brep","step")] = _brep_to_step
    GRAPH[("brep","stl")]  = _brep_to_stl
    GRAPH[("iges","step")] = _iges_to_step
    GRAPH[("stl","step")]  = _stl_to_step
except Exception: pass

if _EZDXF:
    GRAPH[("dxf","svg")] = _dxf_to_svg
    GRAPH[("dxf","txt")] = _dxf_to_txt
    try:
        import matplotlib
        GRAPH[("dxf","png")] = _dxf_to_png
        GRAPH[("dxf","pdf")] = _dxf_to_pdf
    except ImportError:
        pass
    if _TRIMESH:
        for _b in ["stl","obj","ply","glb","3mf","off"]:
            GRAPH[("dxf",_b)] = _dxf_to_mesh
        for _a in _MESH_FMTS:
            GRAPH[(_a,"dxf")] = _mesh_to_dxf

if _INKSCAPE:
    GRAPH[("svg","dxf")] = _svg_to_dxf

if _IFCOS and _TRIMESH:
    for b in ["obj","stl","ply","glb"]:
        GRAPH[("ifc",b)] = _ifc_to_obj
    GRAPH[("ifc","glb")] = _ifc_to_glb
    GRAPH[("ifczip","glb")] = _ifc_to_glb

# ── Blender CLI ───────────────────────────────────────────────────────────────
if _BLENDER:
    for a in ["fbx","blend","dae","glb"]:
        for b in ["glb","fbx","stl","obj","dae","usdz","ply","x3d"]:
            if a != b:
                GRAPH[(a,b)] = _blend_export
    GRAPH[("blend","abc")] = _blend_export

# ── GIS ───────────────────────────────────────────────────────────────────────
if _PYSHP:
    GRAPH[("shp","geojson")] = _shp_to_geojson
    if _PANDAS:
        def _shp_to_csv(src, dst, **kw):
            tmp = dst + "_tmp.geojson"
            try: _shp_to_geojson(src, tmp); _data_convert(tmp, dst, **kw)
            finally:
                try: os.unlink(tmp)
                except OSError: pass
        GRAPH[("shp","csv")] = _shp_to_csv

GRAPH[("geojson","kml")] = _geojson_to_kml
GRAPH[("kml","geojson")] = _kml_to_geojson

if _FIONA:
    GRAPH[("geojson","shp")] = _geojson_to_shp_fiona
elif _PYSHP:
    GRAPH[("geojson","shp")] = _geojson_to_shp

if _GPXPY:
    GRAPH[("gpx","geojson")] = _gpx_to_geojson
    GRAPH[("gpx","csv")]     = _gpx_to_csv

GRAPH[("kmz","kml")] = lambda s,d,**k: (
    __import__("zipfile").ZipFile(s).extract(
        [n for n in __import__("zipfile").ZipFile(s).namelist() if n.endswith(".kml")][0],
        __import__("pathlib").Path(d).parent
    )
)

# ── Fonts ─────────────────────────────────────────────────────────────────────
if _FONTTOOLS:
    for a in _FONT_FMTS:
        for b in _FONT_FMTS:
            if a != b: GRAPH[(a,b)] = _font_convert

# ── Stickerei ─────────────────────────────────────────────────────────────────
if _PYEMB:
    for a in _EMBD_FMTS:
        for b in _EMBD_FMTS:
            if a != b: GRAPH[(a,b)] = _embroidery_convert
        GRAPH[(a,"svg")] = _embroidery_to_svg
        GRAPH[(a,"csv")] = _embroidery_to_csv

# ── Medizin ───────────────────────────────────────────────────────────────────
if _PYDICOM and _PIL:
    GRAPH[("dcm","png")]  = _dcm_to_png
    GRAPH[("dcm","jpeg")] = _dcm_to_png
    GRAPH[("dcm","tiff")] = _dcm_to_png

# ── LiDAR ─────────────────────────────────────────────────────────────────────
if _LASPY:
    GRAPH[("las","csv")] = _las_to_csv
    GRAPH[("laz","las")] = _laz_to_las
    if _TRIMESH:
        GRAPH[("las","ply")] = _las_to_ply

# ── PCB / Gerber ──────────────────────────────────────────────────────────────
if _GERBER:
    for a in _GERBER_FMTS:
        GRAPH[(a,"png")] = _gerber_to_png
        GRAPH[(a,"svg")] = _gerber_to_svg

# ── G-Code ───────────────────────────────────────────────────────────────────
for a in ["gcode","nc","cnc","tap","ngc","mpf"]:
    GRAPH[(a,"svg")] = _gcode_to_svg
    GRAPH[(a,"csv")] = _gcode_to_csv

# ── Archive ───────────────────────────────────────────────────────────────────
for a in _IMG_FMTS + _AUDIO_FMTS + _DATA_FMTS + ["pdf","docx","xlsx","pptx","svg"]:
    GRAPH[(a,"zip")] = _to_zip

if _PY7ZR:
    for a in _IMG_FMTS + _AUDIO_FMTS + _DATA_FMTS + ["pdf","docx","xlsx"]:
        GRAPH[(a,"7z")] = _to_7z
    GRAPH[("7z","zip")] = lambda s,d,**k: _7z_extract_first(s,d)

if _RARFILE:
    GRAPH[("rar","zip")] = lambda s,d,**k: _rar_extract_first(s,d)

GRAPH[("zip","7z")]     = _to_7z
GRAPH[("zip","tar")]    = _zip_to_tar
GRAPH[("tar","7z")]     = _to_7z
GRAPH[("tar_gz","7z")]  = _to_7z
GRAPH[("tar_bz2","7z")] = _to_7z

# ── PDF → alle Seiten als ZIP ─────────────────────────────────────────────────
if _PYMUPDF:
    def _pdf_to_pages_zip(src, dst, **kw):
        import fitz, zipfile
        doc = fitz.open(src)
        n = len(doc)
        pad = len(str(n))
        with zipfile.ZipFile(dst, "w", zipfile.ZIP_DEFLATED) as zf:
            for i, page in enumerate(doc):
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                zf.writestr(f"seite_{str(i+1).zfill(pad)}.png", pix.tobytes("png"))
    GRAPH[("pdf","zip")] = _pdf_to_pages_zip  # override generisches to_zip für PDF

# ── Hintergrund entfernen ─────────────────────────────────────────────────────
if _REMBG and _PIL:
    def _remove_bg(src, dst, **kw):
        from rembg import remove
        from PIL import Image
        import io
        with open(src, "rb") as f:
            data = f.read()
        result = remove(data)
        img = Image.open(io.BytesIO(result)).convert("RGBA")
        img.save(dst, "PNG")
    for a in _IMG_FMTS + ["heic","heif","avif"] + _RAW_FMTS:
        GRAPH[(a, "png_nobg")] = _remove_bg

# ── Untertitel ────────────────────────────────────────────────────────────────
_SUB_FMTS = ["srt", "vtt", "ass", "ssa", "sub"]

def _srt_to_vtt(src, dst, **kw):
    import re
    t = Path(src).read_text(encoding="utf-8", errors="replace")
    t = re.sub(r"(\d{2}:\d{2}:\d{2}),(\d{3})", r"\1.\2", t)
    Path(dst).write_text("WEBVTT\n\n" + t.strip() + "\n", encoding="utf-8")

def _vtt_to_srt(src, dst, **kw):
    import re
    t = Path(src).read_text(encoding="utf-8", errors="replace")
    t = re.sub(r"^WEBVTT[^\n]*\n", "", t, flags=re.MULTILINE)
    t = re.sub(r"(\d{2}:\d{2}:\d{2})\.(\d{3})", r"\1,\2", t)
    blocks = [b.strip() for b in t.split("\n\n") if "-->" in b]
    out = []
    for i, b in enumerate(blocks, 1):
        lines = b.split("\n")
        # Entferne WebVTT-Cue-IDs (nicht Timestamps)
        if lines and "-->" not in lines[0]:
            lines = lines[1:]
        out.append(f"{i}\n" + "\n".join(lines))
    Path(dst).write_text("\n\n".join(out) + "\n", encoding="utf-8")

def _srt_to_ass(src, dst, **kw):
    import re
    t = Path(src).read_text(encoding="utf-8", errors="replace")
    header = ("[Script Info]\nScriptType: v4.00+\nPlayResX: 1920\nPlayResY: 1080\n\n"
              "[V4+ Styles]\nFormat: Name, Fontname, Fontsize, PrimaryColour, SecondaryColour, "
              "OutlineColour, BackColour, Bold, Italic, Underline, StrikeOut, ScaleX, ScaleY, "
              "Spacing, Angle, BorderStyle, Outline, Shadow, Alignment, MarginL, MarginR, MarginV, Encoding\n"
              "Style: Default,Arial,48,&H00FFFFFF,&H000000FF,&H00000000,&H64000000,"
              "-1,0,0,0,100,100,0,0,1,2,2,2,10,10,10,1\n\n"
              "[Events]\nFormat: Layer, Start, End, Style, Name, MarginL, MarginR, MarginV, Effect, Text\n")
    def ts(s):
        s = s.replace(",", ".")
        h, m, rest = s.split(":")
        sec, cs = rest.split(".")
        return f"{int(h)}:{m}:{sec}.{cs[:2]}"
    blocks = re.findall(r"\d+\n(\d{2}:\d{2}:\d{2},\d{3}) --> (\d{2}:\d{2}:\d{2},\d{3})\n(.*?)(?=\n\d+\n|\Z)",
                        t, re.DOTALL)
    def clean(txt):
        txt = re.sub("<[^>]+>", "", txt)
        return txt.replace(chr(10), r"{\N}")
    events = [f"Dialogue: 0,{ts(s)},{ts(e)},Default,,0,0,0,,{clean(txt)}"
              for s, e, txt in blocks]
    Path(dst).write_text(header + "\n".join(events) + "\n", encoding="utf-8")

def _ass_to_srt(src, dst, **kw):
    import re
    t = Path(src).read_text(encoding="utf-8", errors="replace")
    def ts(s):
        h, m, rest = s.split(":")
        sec, cs = rest.split(".")
        return f"{int(h):02d}:{m}:{sec},{cs}0"
    events = re.findall(r"Dialogue:.*?,(\d+:\d{2}:\d{2}\.\d{2}),(\d+:\d{2}:\d{2}\.\d{2}),.*?,.*?,.*?,.*?,.*?,.*?,(.*)", t)
    out = []
    for i, (s, e, txt) in enumerate(events, 1):
        clean = re.sub(r"\{[^}]+\}", "", txt).replace(r"{\N}", "\n")
        out.append(f"{i}\n{ts(s)} --> {ts(e)}\n{clean}")
    Path(dst).write_text("\n\n".join(out) + "\n", encoding="utf-8")

def _sub_to_srt(src, dst, **kw):
    import re
    t = Path(src).read_text(encoding="utf-8", errors="replace")
    lines = t.splitlines()
    out, counter = [], 1
    fps = 25.0
    for line in lines:
        m = re.match(r"\{(\d+)\}\{(\d+)\}(.*)", line)
        if m:
            f1, f2, text = int(m.group(1)), int(m.group(2)), m.group(3)
            def fr(f):
                s = f / fps
                h, r = divmod(s, 3600)
                m2, s2 = divmod(r, 60)
                return f"{int(h):02d}:{int(m2):02d}:{s2:06.3f}".replace(".", ",")
            out.append(f"{counter}\n{fr(f1)} --> {fr(f2)}\n{text.replace('|', chr(10))}\n")
            counter += 1
    Path(dst).write_text("\n".join(out), encoding="utf-8")

GRAPH[("srt","vtt")] = _srt_to_vtt
GRAPH[("vtt","srt")] = _vtt_to_srt
GRAPH[("srt","ass")] = _srt_to_ass
GRAPH[("srt","ssa")] = _srt_to_ass
GRAPH[("ass","srt")] = _ass_to_srt
GRAPH[("ssa","srt")] = _ass_to_srt
GRAPH[("sub","srt")] = _sub_to_srt
def _vtt_to_ass(src, dst, **kw):
    tmp = src + "_tmp.srt"
    try: _vtt_to_srt(src, tmp); _srt_to_ass(tmp, dst)
    finally:
        try: os.unlink(tmp)
        except OSError: pass
GRAPH[("vtt","ass")] = _vtt_to_ass

# Untertitel → TXT (plain text)
def _sub_to_txt(src, dst, **kw):
    import re
    t = Path(src).read_text(encoding="utf-8", errors="replace")
    # Entferne Timestamps und Sequenznummern
    t = re.sub(r"\d+\n\d{2}:\d{2}:\d{2}[,\.]\d{3} --> \d{2}:\d{2}:\d{2}[,\.]\d{3}\n", "", t)
    t = re.sub(r"WEBVTT[^\n]*\n", "", t)
    t = re.sub(r"\{[^}]+\}", "", t)
    t = re.sub(r"<[^>]+>", "", t)
    Path(dst).write_text(t.strip(), encoding="utf-8")

for a in _SUB_FMTS:
    GRAPH[(a,"txt")] = _sub_to_txt

# ── Chemie-/Molekülformate ────────────────────────────────────────────────────
def _mol_to_json(src, dst, **kw):
    """MOL/SDF → JSON (Atome, Bindungen) — kein rdkit erforderlich"""
    t = Path(src).read_text(encoding="utf-8", errors="replace")
    blocks = t.split("$$$$")
    molecules = []
    for block in blocks:
        lines = block.strip().splitlines()
        if len(lines) < 4: continue
        name = lines[0].strip()
        try:
            counts = lines[3].split()
            n_atoms = int(counts[0])
            n_bonds = int(counts[1]) if len(counts) > 1 else 0
        except (ValueError, IndexError): continue
        atoms, bonds = [], []
        for l in lines[4:4+n_atoms]:
            parts = l.split()
            if len(parts) >= 4:
                atoms.append({"x": float(parts[0]), "y": float(parts[1]),
                               "z": float(parts[2]), "element": parts[3]})
        for l in lines[4+n_atoms:4+n_atoms+n_bonds]:
            parts = l.split()
            if len(parts) >= 3:
                bonds.append({"from": int(parts[0])-1, "to": int(parts[1])-1,
                               "type": int(parts[2])})
        molecules.append({"name": name, "atoms": atoms, "bonds": bonds})
    Path(dst).write_text(json.dumps(molecules if len(molecules)>1 else (molecules[0] if molecules else {}),
                                   indent=2), encoding="utf-8")

def _pdb_to_json(src, dst, **kw):
    """PDB (Protein Data Bank) → JSON"""
    atoms = []
    for line in Path(src).read_text(encoding="utf-8", errors="replace").splitlines():
        if line.startswith(("ATOM  ", "HETATM")):
            try:
                atoms.append({
                    "type": line[:6].strip(),
                    "serial": int(line[6:11]),
                    "name": line[12:16].strip(),
                    "resName": line[17:20].strip(),
                    "chainID": line[21].strip(),
                    "resSeq": int(line[22:26]),
                    "x": float(line[30:38]),
                    "y": float(line[38:46]),
                    "z": float(line[46:54]),
                    "element": line[76:78].strip() if len(line) > 76 else "",
                })
            except (ValueError, IndexError): pass
    Path(dst).write_text(json.dumps({"atoms": atoms, "count": len(atoms)}, indent=2),
                         encoding="utf-8")

GRAPH[("mol","json")] = _mol_to_json
GRAPH[("sdf","json")] = _mol_to_json
GRAPH[("pdb","json")] = _pdb_to_json

# ── Wissenschaftliche Formate ─────────────────────────────────────────────────
if _H5PY:
    def _hdf5_to_json(src, dst, **kw):
        import h5py, json
        def _serialize(obj):
            if hasattr(obj, "items"):
                return {k: _serialize(v) for k, v in obj.items()}
            elif hasattr(obj, "__len__") and not isinstance(obj, str):
                arr = obj[:]
                return arr.tolist() if hasattr(arr, "tolist") else list(arr)
            return obj
        with h5py.File(src, "r") as f:
            data = _serialize(f)
        Path(dst).write_text(json.dumps(data, indent=2, default=str), encoding="utf-8")
    GRAPH[("h5","json")]   = _hdf5_to_json
    GRAPH[("hdf5","json")] = _hdf5_to_json
    GRAPH[("nc","json")]   = _hdf5_to_json  # NetCDF4 ist HDF5-kompatibel

# ══════════ YAML / TOML / CONFIG ══════════

if _YAML:
    def _yaml_to_json(src, dst, **kw):
        import yaml
        data = yaml.safe_load(Path(src).read_text(encoding="utf-8"))
        Path(dst).write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")

    def _json_to_yaml(src, dst, **kw):
        import yaml
        data = json.loads(Path(src).read_text(encoding="utf-8"))
        Path(dst).write_text(yaml.dump(data, allow_unicode=True, default_flow_style=False,
                                        sort_keys=False), encoding="utf-8")

    def _yaml_to_csv(src, dst, **kw):
        import yaml, csv
        data = yaml.safe_load(Path(src).read_text(encoding="utf-8"))
        rows = data if isinstance(data, list) else [data]
        if not rows: raise ValueError("YAML enthält keine Liste")
        with open(dst, "w", newline="", encoding="utf-8") as f:
            w = csv.DictWriter(f, fieldnames=list(rows[0].keys()))
            w.writeheader(); w.writerows(rows)

    GRAPH[("yaml","json")]  = _yaml_to_json
    GRAPH[("yml","json")]   = _yaml_to_json
    GRAPH[("json","yaml")]  = _json_to_yaml
    GRAPH[("yaml","csv")]   = _yaml_to_csv
    if _PANDAS:
        def _yaml_to_data(src, dst, **kw):
            tmp = src + "_tmp.json"
            try: _yaml_to_json(src, tmp); _data_convert(tmp, dst, **kw)
            finally:
                try: os.unlink(tmp)
                except OSError: pass
        def _csv_to_yaml(src, dst, **kw):
            tmp = src + "_tmp.json"
            try: _data_convert(src, tmp, **kw); _json_to_yaml(tmp, dst, **kw)
            finally:
                try: os.unlink(tmp)
                except OSError: pass
        for d in ("csv","tsv","xlsx","xml","parquet"):
            GRAPH[("yaml",d)] = _yaml_to_data
        GRAPH[("csv","yaml")] = _csv_to_yaml

# TOML ↔ JSON — tomllib ist Python 3.11+ stdlib
def _toml_to_json(src, dst, **kw):
    import tomllib
    data = tomllib.loads(Path(src).read_text(encoding="utf-8"))
    Path(dst).write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")

def _toml_to_yaml(src, dst, **kw):
    import tomllib, yaml
    data = tomllib.loads(Path(src).read_text(encoding="utf-8"))
    Path(dst).write_text(yaml.dump(data, allow_unicode=True, default_flow_style=False), encoding="utf-8")

def _json_to_toml(src, dst, **kw):
    try:
        import tomli_w
        data = json.loads(Path(src).read_text(encoding="utf-8"))
        if isinstance(data, list):
            data = {"data": data}
        content = tomli_w.dumps(data)
        Path(dst).write_text(content or "# keine Daten\n", encoding="utf-8")
    except ImportError:
        raise RuntimeError("pip install tomli-w  für JSON→TOML")

GRAPH[("toml","json")] = _toml_to_json
GRAPH[("json","toml")] = _json_to_toml
if _YAML:
    GRAPH[("toml","yaml")] = _toml_to_yaml
    GRAPH[("yaml","toml")] = lambda s,o,**k: (_yaml_to_json(s,s+"_tmp.json") or _json_to_toml(s+"_tmp.json",o))

# ══════════ SVG → PDF ══════════

def _svg_to_pdf(src, dst, **kw):
    """SVG → PDF via Inkscape oder cairosvg"""
    if _INKSCAPE:
        r = _run(["inkscape", "--export-type=pdf", f"--export-filename={dst}", src])
        if r.returncode == 0 and Path(dst).exists(): return
    try:
        import cairosvg
        cairosvg.svg2pdf(url=src, write_to=dst); return
    except ImportError:
        pass
    raise RuntimeError("SVG→PDF: inkscape oder cairosvg erforderlich")

GRAPH[("svg","pdf")] = _svg_to_pdf

# ══════════ VSDX (Visio) ══════════

def _vsdx_to_json(src, dst, **kw):
    """Visio VSDX → JSON (Shapes, Text, Verbindungen)"""
    import vsdx
    result = {"pages": []}
    with vsdx.VisioFile(src) as vis:
        for page in vis.pages:
            shapes_data = []
            for shape in page.shapes:
                shapes_data.append({
                    "id":   shape.ID,
                    "text": (shape.text or "").strip(),
                    "x":    shape.x,
                    "y":    shape.y,
                    "label": (getattr(shape, "label", "") or "").strip(),
                })
            result["pages"].append({"name": page.name, "shapes": shapes_data})
    Path(dst).write_text(json.dumps(result, indent=2, default=str, ensure_ascii=False), encoding="utf-8")

def _vsdx_to_txt(src, dst, **kw):
    """Visio VSDX → TXT (alle Shape-Texte)"""
    import vsdx
    texts = []
    with vsdx.VisioFile(src) as vis:
        for page in vis.pages:
            texts.append(f"=== {page.name} ===")
            for shape in page.shapes:
                t = (shape.text or "").strip()
                if t: texts.append(t)
    Path(dst).write_text("\n".join(texts), encoding="utf-8")

if _VSDX:
    GRAPH[("vsdx","json")] = _vsdx_to_json
    GRAPH[("vsdx","txt")]  = _vsdx_to_txt
    if _SOFFICE:
        GRAPH[("vsdx","pdf")] = _libreoffice
        GRAPH[("vsdx","svg")] = _libreoffice

# ══════════ MSG / EML (E-Mail) ══════════

def _msg_to_txt(src, dst, **kw):
    """Outlook MSG → TXT via extract-msg"""
    import extract_msg
    msg = extract_msg.openMsg(src)
    lines = []
    if msg.subject:  lines.append(f"Betreff: {msg.subject}")
    if msg.sender:   lines.append(f"Von:     {msg.sender}")
    if msg.to:       lines.append(f"An:      {msg.to}")
    if msg.cc:       lines.append(f"CC:      {msg.cc}")
    if msg.date:     lines.append(f"Datum:   {msg.date}")
    lines.append("")
    body = msg.body or ""
    lines.append(body)
    Path(dst).write_text("\n".join(lines), encoding="utf-8")

def _msg_to_html(src, dst, **kw):
    """Outlook MSG → HTML"""
    import extract_msg, html as html_mod
    msg = extract_msg.openMsg(src)
    def esc(s): return html_mod.escape(str(s or ""))
    body_html = (getattr(msg, "htmlBody", None) or "").strip()
    if not body_html:
        body_html = "<pre>" + esc(msg.body or "") + "</pre>"
    content = f"""<!DOCTYPE html><html><head><meta charset="utf-8"><title>{esc(msg.subject)}</title>
<style>body{{font-family:sans-serif;max-width:860px;margin:2rem auto}}
table{{border-collapse:collapse;width:100%;margin-bottom:1rem}}
td,th{{padding:6px 12px;border:1px solid #ddd}}th{{background:#f5f5f5;text-align:left}}</style>
</head><body>
<table>
<tr><th>Betreff</th><td>{esc(msg.subject)}</td></tr>
<tr><th>Von</th><td>{esc(msg.sender)}</td></tr>
<tr><th>An</th><td>{esc(msg.to)}</td></tr>
<tr><th>Datum</th><td>{esc(msg.date)}</td></tr>
</table>
{body_html}
</body></html>"""
    Path(dst).write_text(content, encoding="utf-8")

def _eml_to_txt(src, dst, **kw):
    """EML → TXT (Standard-E-Mail-Format)"""
    import email
    msg = email.message_from_bytes(Path(src).read_bytes())
    lines = [
        f"Betreff: {msg.get('Subject','')}",
        f"Von:     {msg.get('From','')}",
        f"An:      {msg.get('To','')}",
        f"Datum:   {msg.get('Date','')}",
        "",
    ]
    for part in msg.walk():
        if part.get_content_type() == "text/plain":
            payload = part.get_payload(decode=True)
            if payload:
                charset = part.get_content_charset() or "utf-8"
                lines.append(payload.decode(charset, errors="replace"))
    Path(dst).write_text("\n".join(lines), encoding="utf-8")

def _eml_to_html(src, dst, **kw):
    """EML → HTML"""
    import email, html as html_mod
    msg = email.message_from_bytes(Path(src).read_bytes())
    def esc(s): return html_mod.escape(str(s or ""))
    body_html, body_txt = "", ""
    for part in msg.walk():
        ct = part.get_content_type()
        payload = part.get_payload(decode=True)
        if not payload: continue
        charset = part.get_content_charset() or "utf-8"
        decoded = payload.decode(charset, errors="replace")
        if ct == "text/html": body_html = decoded
        elif ct == "text/plain": body_txt = decoded
    body = body_html or ("<pre>" + esc(body_txt) + "</pre>")
    Path(dst).write_text(f"""<!DOCTYPE html><html><head><meta charset="utf-8">
<title>{esc(msg.get("Subject",""))}</title></head><body>
<table border="1" cellpadding="6" style="border-collapse:collapse;margin-bottom:1rem">
<tr><th>Betreff</th><td>{esc(msg.get("Subject",""))}</td></tr>
<tr><th>Von</th><td>{esc(msg.get("From",""))}</td></tr>
<tr><th>An</th><td>{esc(msg.get("To",""))}</td></tr>
<tr><th>Datum</th><td>{esc(msg.get("Date",""))}</td></tr>
</table>{body}</body></html>""", encoding="utf-8")

if _EXTRACTMSG:
    GRAPH[("msg","txt")]  = _msg_to_txt
    GRAPH[("msg","html")] = _msg_to_html
    if _PANDOC or _REPORTLAB:
        def _msg_to_pdf(src, dst, **kw):
            tmp = src + "_tmp.html"
            try:
                _msg_to_html(src, tmp)
                _pandoc(tmp, dst) if _PANDOC else _html_to_pdf(tmp, dst)
            finally:
                try: os.unlink(tmp)
                except OSError: pass
        GRAPH[("msg","pdf")] = _msg_to_pdf

GRAPH[("eml","txt")]  = _eml_to_txt
GRAPH[("eml","html")] = _eml_to_html

# ══════════ TAR-ARCHIVE ══════════

def _to_tar(src, dst, **kw):
    import tarfile
    suf = Path(dst).suffixes
    mode = "w:gz" if ".gz" in suf or dst.endswith(".tgz") else \
           "w:bz2" if ".bz2" in suf or dst.endswith(".tbz2") else \
           "w:xz" if ".xz" in suf or dst.endswith(".txz") else "w"
    with tarfile.open(dst, mode) as t:
        t.add(src, arcname=Path(src).name)

def _tar_extract_first(src, dst, **kw):
    import tarfile
    with tarfile.open(src) as t:
        names = [m.name for m in t.getmembers() if m.isfile()]
        if not names: raise ValueError("TAR ist leer")
        with tempfile.TemporaryDirectory() as td:
            t.extractall(td, filter="data")
            shutil.copy(os.path.join(td, names[0]), dst)

_TAR_FMTS = ["tar", "tar_gz", "tar_bz2", "tar_xz"]
for a in _IMG_FMTS + _AUDIO_FMTS + _DATA_FMTS + ["pdf","docx","svg"]:
    GRAPH[(a,"tar")] = _to_tar
GRAPH[("tar","zip")]     = _tar_extract_first
GRAPH[("tar_gz","zip")]  = _tar_extract_first
GRAPH[("tar_bz2","zip")] = _tar_extract_first

# ══════════ VCF (vCard — Kontakte) ══════════

def _vcf_to_json(src, dst, **kw):
    """vCard VCF → JSON"""
    text = Path(src).read_text(encoding="utf-8", errors="replace")
    contacts, current = [], {}
    for line in text.splitlines():
        line = line.strip()
        if line == "BEGIN:VCARD":
            current = {}
        elif line == "END:VCARD":
            if current: contacts.append(current)
        elif ":" in line:
            key, _, val = line.partition(":")
            base = key.split(";")[0].upper()
            val = val.strip()
            if base in ("FN","N","EMAIL","TEL","ORG","ADR","BDAY","URL","NOTE","TITLE","NICKNAME"):
                if base in current:
                    if not isinstance(current[base], list): current[base] = [current[base]]
                    current[base].append(val)
                else:
                    current[base] = val
    Path(dst).write_text(json.dumps(contacts, indent=2, ensure_ascii=False), encoding="utf-8")

def _vcf_to_csv(src, dst, **kw):
    import csv
    tmp = src + "_tmp.json"
    try:
        _vcf_to_json(src, tmp)
        items = json.loads(Path(tmp).read_text())
        if isinstance(items, dict): items = [items]
        if not items: raise ValueError("VCF enthält keine Kontakte")
        keys = list({k for c in items for k in c.keys()})
        with open(dst, "w", newline="", encoding="utf-8") as f:
            w = csv.DictWriter(f, fieldnames=keys)
            w.writeheader()
            for c in items:
                w.writerow({k: ("; ".join(v) if isinstance(v,list) else v) for k,v in c.items()})
    finally:
        try: os.unlink(tmp)
        except OSError: pass

GRAPH[("vcf","json")] = _vcf_to_json
GRAPH[("vcf","csv")]  = _vcf_to_csv
GRAPH[("vcf","txt")]  = lambda s,d,**k: Path(d).write_text(Path(s).read_text(encoding="utf-8",errors="replace"), encoding="utf-8")

# ══════════ ICS (iCalendar — Kalender) ══════════

def _ics_to_json(src, dst, **kw):
    """ICS → JSON (Events, To-Dos)"""
    try:
        from icalendar import Calendar
        cal = Calendar.from_ical(Path(src).read_bytes())
        events = []
        for comp in cal.walk():
            if comp.name in ("VEVENT","VTODO","VJOURNAL"):
                def gs(field):
                    v = comp.get(field)
                    if v is None: return ""
                    if hasattr(v, "dt"): return str(v.dt)
                    return str(v)
                events.append({
                    "type":        comp.name,
                    "summary":     gs("SUMMARY"),
                    "start":       gs("DTSTART"),
                    "end":         gs("DTEND"),
                    "location":    gs("LOCATION"),
                    "description": gs("DESCRIPTION"),
                    "uid":         gs("UID"),
                    "status":      gs("STATUS"),
                })
        Path(dst).write_text(json.dumps(events, indent=2, ensure_ascii=False), encoding="utf-8")
    except ImportError:
        # Pure-Python fallback: basic line parsing
        import re
        events, current = [], {}
        for line in Path(src).read_text(encoding="utf-8",errors="replace").splitlines():
            if line.startswith("BEGIN:VEVENT"): current = {"type":"VEVENT"}
            elif line.startswith("END:VEVENT"):
                if current: events.append(current)
            elif ":" in line:
                k, _, v = line.partition(":")
                current[k.split(";")[0]] = v
        Path(dst).write_text(json.dumps(events, indent=2, ensure_ascii=False), encoding="utf-8")

def _ics_to_csv(src, dst, **kw):
    import csv
    tmp = src + "_tmp.json"
    try:
        _ics_to_json(src, tmp)
        events = json.loads(Path(tmp).read_text())
        if not events: raise ValueError("ICS enthält keine Einträge")
        keys = list(events[0].keys())
        with open(dst, "w", newline="", encoding="utf-8") as f:
            w = csv.DictWriter(f, fieldnames=keys)
            w.writeheader(); w.writerows(events)
    finally:
        try: os.unlink(tmp)
        except OSError: pass

GRAPH[("ics","json")] = _ics_to_json
GRAPH[("ics","csv")]  = _ics_to_csv
GRAPH[("ics","txt")]  = lambda s,d,**k: Path(d).write_text(Path(s).read_text(encoding="utf-8",errors="replace"), encoding="utf-8")

# ══════════ FB2 (FictionBook E-Book) ══════════

def _fb2_to_txt(src, dst, **kw):
    """FictionBook2 → TXT (alle Absätze)"""
    import re
    xml = Path(src).read_text(encoding="utf-8", errors="replace")
    text = re.sub(r"<title>(.*?)</title>", r"\n\n== \1 ==\n", xml, flags=re.DOTALL)
    text = re.sub(r"<[^>]+>", "", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    Path(dst).write_text(text.strip(), encoding="utf-8")

def _fb2_to_html(src, dst, **kw):
    """FictionBook2 → HTML"""
    import re
    xml = Path(src).read_text(encoding="utf-8", errors="replace")
    html = re.sub(r"<\?xml[^>]*\?>", "", xml)
    html = re.sub(r"<FictionBook[^>]*>", "<!DOCTYPE html><html><head><meta charset='utf-8'></head><body>", html)
    html = html.replace("</FictionBook>", "</body></html>")
    html = html.replace("<title>", "<h2>").replace("</title>", "</h2>")
    html = html.replace("<section>", "<section style='margin:1rem 0'>")
    html = html.replace("<subtitle>", "<h3>").replace("</subtitle>", "</h3>")
    html = html.replace("<strong>", "<strong>").replace("<emphasis>", "<em>").replace("</emphasis>", "</em>")
    html = re.sub(r"<(image|binary)[^>]*/?>", "", html)
    Path(dst).write_text(html, encoding="utf-8")

GRAPH[("fb2","txt")]  = _fb2_to_txt
GRAPH[("fb2","html")] = _fb2_to_html
if _PANDOC:
    GRAPH[("fb2","epub")]     = _pandoc
    GRAPH[("fb2","markdown")] = _pandoc

# ══════════ QR-CODE ══════════

def _txt_to_qr(src, dst, **kw):
    """TXT/URL → QR-Code PNG via qrcode"""
    try:
        import qrcode
        text = Path(src).read_text(encoding="utf-8", errors="replace").strip()
        # QR v40 ERROR_CORRECT_L: max 2953 bytes binary. Truncate on byte level to be safe.
        MAX_BYTES = 2900
        encoded = text.encode("utf-8")
        if len(encoded) > MAX_BYTES:
            text = encoded[:MAX_BYTES].decode("utf-8", errors="ignore") + "..."
        qr = qrcode.QRCode(
            version=None, error_correction=qrcode.constants.ERROR_CORRECT_L,
            box_size=10, border=4,
        )
        qr.add_data(text)
        qr.make(fit=True)
        img = qr.make_image(fill_color="black", back_color="white")
        img.save(dst)
    except ImportError:
        raise RuntimeError("pip install qrcode[pil]")

if _QRCODE and _PIL:
    GRAPH[("txt","qr_png")] = _txt_to_qr
    GRAPH[("html","qr_png")] = _txt_to_qr
    GRAPH[("url","qr_png")]  = _txt_to_qr

# ══════════ DXF ERWEITERUNG (BIM/CAD) ══════════

def _dxf_to_json(src, dst, **kw):
    """DXF → JSON (alle Entitäten mit Koordinaten und Attributen)"""
    import ezdxf
    doc = ezdxf.readfile(src)
    entities = []
    for entity in doc.modelspace():
        e = {"type": entity.dxftype()}
        try:
            if hasattr(entity.dxf, "start"):
                e["start"] = list(entity.dxf.start)
            if hasattr(entity.dxf, "end"):
                e["end"] = list(entity.dxf.end)
            if hasattr(entity.dxf, "center"):
                e["center"] = list(entity.dxf.center)
            if hasattr(entity.dxf, "radius"):
                e["radius"] = entity.dxf.radius
            if hasattr(entity.dxf, "text"):
                e["text"] = entity.dxf.text
            if hasattr(entity.dxf, "layer"):
                e["layer"] = entity.dxf.layer
            if hasattr(entity.dxf, "color"):
                e["color"] = entity.dxf.color
        except Exception:
            pass
        entities.append(e)
    Path(dst).write_text(json.dumps({"entities": entities, "count": len(entities)},
                                    indent=2, default=str, ensure_ascii=False), encoding="utf-8")

if _EZDXF:
    GRAPH[("dxf","json")] = _dxf_to_json
    # dxf→pdf und dxf→png bereits via Matplotlib oben registriert — kein Override

# ══════════ G-CODE ERWEITERUNG ══════════

def _gcode_to_json(src, dst, **kw):
    """G-Code → JSON (vollständiger Parse: Koordinaten, Feedrate, Spindle, Kommentare)"""
    import re
    lines = Path(src).read_text(encoding="utf-8", errors="replace").split("\n")
    moves, x, y, z, f = [], 0.0, 0.0, 0.0, 0.0
    tool, spindle = 0, 0
    for raw in lines:
        comment_match = re.search(r"[;(](.+)", raw)
        comment = comment_match.group(1).strip().rstrip(")") if comment_match else ""
        line = re.sub(r"[;(].*", "", raw).strip().upper()
        if not line: continue
        # Tool change
        tm = re.search(r"T(\d+)", line)
        if tm: tool = int(tm.group(1))
        sm = re.search(r"S(\d+)", line)
        if sm: spindle = int(sm.group(1))
        fm = re.search(r"F([\d.]+)", line)
        if fm: f = float(fm.group(1))
        if re.match(r"G[01][^0-9]|^G[01]$", line) or re.search(r"[XYZ][-\d.]", line):
            prev_x, prev_y, prev_z = x, y, z
            for ax, var in [("X",x),("Y",y),("Z",z)]:
                m = re.search(rf"{ax}([-\d.]+)", line)
                if m:
                    val = float(m.group(1))
                    if ax == "X": x = val
                    elif ax == "Y": y = val
                    elif ax == "Z": z = val
            mode = "rapid" if re.match(r"G0[^1]|^G0$", line) else "cut"
            if x != prev_x or y != prev_y or z != prev_z:
                moves.append({"mode": mode, "x": x, "y": y, "z": z,
                               "feedrate": f, "tool": tool, "spindle": spindle,
                               "comment": comment})
    Path(dst).write_text(json.dumps({"moves": moves, "count": len(moves)}, indent=2), encoding="utf-8")

for a in ["gcode","nc","cnc","tap","ngc","mpf"]:
    GRAPH[(a,"json")] = _gcode_to_json

# ══════════ GIS ERWEITERUNG ══════════

def _gpx_to_kml(src, dst, **kw):
    """GPX → KML"""
    import gpxpy
    gpx = gpxpy.parse(Path(src).read_text(encoding="utf-8"))
    lines = ['<?xml version="1.0" encoding="UTF-8"?>',
             '<kml xmlns="http://www.opengis.net/kml/2.2"><Document>']
    for track in gpx.tracks:
        lines.append(f"<Placemark><name>{track.name or 'Track'}</name>")
        coords = " ".join(f"{p.longitude},{p.latitude},{p.elevation or 0}"
                          for seg in track.segments for p in seg.points)
        lines.append(f"<LineString><coordinates>{coords}</coordinates></LineString></Placemark>")
    for wpt in gpx.waypoints:
        lines.append(f"<Placemark><name>{wpt.name or 'Waypoint'}</name>"
                     f"<Point><coordinates>{wpt.longitude},{wpt.latitude}</coordinates></Point></Placemark>")
    lines.append("</Document></kml>")
    Path(dst).write_text("\n".join(lines), encoding="utf-8")

def _geojson_to_csv(src, dst, **kw):
    """GeoJSON → CSV (Punkt-Features mit Eigenschaften)"""
    import csv
    data = json.loads(Path(src).read_text(encoding="utf-8"))
    features = data.get("features", [])
    if not features: raise ValueError("GeoJSON enthält keine Features")
    rows = []
    for feat in features:
        row = dict(feat.get("properties") or {})
        geo = feat.get("geometry", {})
        if geo.get("type") == "Point":
            coords = geo.get("coordinates", [])
            row["longitude"] = coords[0] if len(coords) > 0 else ""
            row["latitude"]  = coords[1] if len(coords) > 1 else ""
        rows.append(row)
    keys = list({k for r in rows for k in r.keys()})
    with open(dst, "w", newline="", encoding="utf-8") as f:
        w = csv.DictWriter(f, fieldnames=keys)
        w.writeheader(); w.writerows(rows)

def _csv_to_geojson(src, dst, **kw):
    """CSV (mit lat/lon Spalten) → GeoJSON"""
    import csv
    features = []
    with open(src, newline="", encoding="utf-8", errors="replace") as f:
        reader = csv.DictReader(f)
        for row in reader:
            lat_key = next((k for k in row if k.lower() in ("lat","latitude")), None)
            lon_key = next((k for k in row if k.lower() in ("lon","lng","longitude")), None)
            if lat_key and lon_key:
                try:
                    coords = [float(row[lon_key]), float(row[lat_key])]
                    props = {k:v for k,v in row.items() if k not in (lat_key, lon_key)}
                    features.append({"type":"Feature","geometry":{"type":"Point","coordinates":coords},"properties":props})
                except ValueError:
                    pass
    Path(dst).write_text(json.dumps({"type":"FeatureCollection","features":features},indent=2,ensure_ascii=False), encoding="utf-8")

if _GPXPY:
    GRAPH[("gpx","kml")] = _gpx_to_kml
    GRAPH[("gpx","xml")] = _gpx_to_kml  # KML ist valides XML

GRAPH[("geojson","csv")]  = _geojson_to_csv
GRAPH[("csv","geojson")]  = _csv_to_geojson
_shp_writer = _geojson_to_shp_fiona if _FIONA else _geojson_to_shp if _PYSHP else None
if _shp_writer:
    def _csv_to_shp(src, dst, **kw):
        tmp = src + "_tmp.geojson"
        try: _csv_to_geojson(src, tmp); _shp_writer(tmp, dst)
        finally:
            try: os.unlink(tmp)
            except OSError: pass
    GRAPH[("csv","shp")] = _csv_to_shp

# ── BFS ──────────────────────────────────────────────────────────────────────

ALL_FORMATS: list = sorted({fmt for pair in GRAPH for fmt in pair})

# Adjazenz-Dict: O(1)-Lookup statt O(N) Iteration über alle Kanten
_ADJ: dict = {}
for (_ga, _gb) in GRAPH:
    _ADJ.setdefault(_ga, []).append(_gb)

# ── LibreOffice Listener vorwärmen ────────────────────────────────────────────
if _SOFFICE and _LO_PYTHON:
    import threading as _lo_thread
    _lo_thread.Thread(target=_ensure_lo_listener, daemon=True).start()

def find_path(src: str, tgt: str):
    if src == tgt: return []
    queue: deque = deque([(src, [])])
    visited: set = {src}
    while queue:
        cur, path = queue.popleft()
        for b in _ADJ.get(cur, []):
            if b not in visited:
                new_path = path + [(cur, b)]
                if b == tgt: return new_path
                visited.add(b); queue.append((b, new_path))
    return None

def list_targets(src: str, max_hops: int = 2) -> list:
    """Gibt alle erreichbaren Zielformate zurück (bis max_hops Schritte)."""
    reachable: set = set()
    queue: deque = deque([(src, 0)])
    visited: set = {src}
    while queue:
        cur, hops = queue.popleft()
        if hops >= max_hops:
            continue
        for b in _ADJ.get(cur, []):
            reachable.add(b)
            if b not in visited:
                visited.add(b)
                queue.append((b, hops + 1))
    reachable.discard(src)
    return sorted(reachable)

# ── Haupt-Konvertierung ───────────────────────────────────────────────────────

_EXT_MAP = {
    "jpeg":"jpg", "tiff":"tif", "markdown":"md",
    "aac":"aac", "step":"step", "iges":"iges",
    "midi":"mid", "sqlite":"db",
    "tar_gz":"tar.gz", "tar_bz2":"tar.bz2", "tar_xz":"tar.xz",
    "qr_png":"png",
}

def convert(file_path: str, target_fmt: str, **kw) -> dict:
    src_fmt = norm(ext_of(file_path))
    tgt_fmt = norm(target_fmt)
    path = find_path(src_fmt, tgt_fmt)

    if path is None:
        return {"error": f"Keine Konvertierung von '{src_fmt}' nach '{tgt_fmt}' möglich."}

    steps = ([path[0][0]] + [b for _,b in path]) if path else [src_fmt]
    ext   = _EXT_MAP.get(tgt_fmt, tgt_fmt)

    with tempfile.TemporaryDirectory() as tmp:
        current = file_path
        for i, (a, b) in enumerate(path):
            out_ext = _EXT_MAP.get(b, b)
            out = os.path.join(tmp, f"s{i}_{a}_{b}.{out_ext}")
            try:
                GRAPH[(a,b)](current, out, **kw)
            except Exception as e:
                return {"error": str(e), "steps": steps}
            if not os.path.exists(out) or os.path.getsize(out) == 0:
                return {"error": f"Schritt {a}→{b} lieferte kein Ergebnis.", "steps": steps}
            current = out

        size = os.path.getsize(current)
        return {
            "result_b64":  _b64(current),
            "result_fmt":  tgt_fmt,
            "result_ext":  ext,
            "steps":       steps,
            "size_bytes":  size,
        }

# ── Stickerei: öffentliche API-Funktionen ─────────────────────────────────────

# Maschinendatenbank: Marke → Modelle + kompatible Formate
MACHINE_DB = {
    "Brother": {
        "formats": ["pes"],
        "models": ["PE800","PE535","PE550D","SE700","SE1900","PR1050X",
                   "Innov-is NQ3600D","Innov-is XV8550D","Stellaire XJ2","Luminaire XP3"]
    },
    "Janome": {
        "formats": ["jef","sew"],
        "models": ["MC500E","MC550E","Skyline S9","Continental M7",
                   "Horizon MC12000","Memory Craft 9850","MC15000"]
    },
    "Pfaff": {
        "formats": ["vp3","vip"],
        "models": ["Creative 4","Creative 4.5","Creative Icon 2",
                   "Performance 5","Quilt Expression 720","Creative Vision"]
    },
    "Husqvarna Viking": {
        "formats": ["vp3","hus"],
        "models": ["Designer Ruby","Designer EPIC 2","Brilliance 80","Sapphire 960Q",
                   "Emerald 158","Designer SE","Topaz 40"]
    },
    "Bernina": {
        "formats": ["exp"],
        "models": ["770 QE Plus","790 Plus","B 880 Plus","535","740","475 QE","L 890"]
    },
    "Elna": {
        "formats": ["jef","exp"],
        "models": ["eXpressive 860+","Excellence 680+","Destiny 880 Plus","Envolve 8300D"]
    },
    "Singer": {
        "formats": ["xxx"],
        "models": ["Futura XL-550","SE300","SE9180Pro","Quantum Stylist 9960"]
    },
    "Baby Lock": {
        "formats": ["pes"],
        "models": ["Solaris XE","Altair","Meridian","Flourish II","Verve","Crown"]
    },
    "Tajima (Industrie)": {
        "formats": ["dst"],
        "models": ["TMEF-HC1501","TCMX-C Series","TFMX","TMAR-KC"]
    },
    "ZSK (Industrie)": {
        "formats": ["dst","dsb"],
        "models": ["Sprint 6","Racer 6","Master","VARIOFRAME"]
    },
    "Toyota": {
        "formats": ["exp"],
        "models": ["ES8460","Oekaki Renaissance","AD-850"]
    },
    "Bernette": {
        "formats": ["exp"],
        "models": ["B70 DECO","B79","B77"]
    },
    "Necchi": {
        "formats": ["exp","pes"],
        "models": ["NA7234","NQ3600D"]
    },
    "Universal (alle)": {
        "formats": ["dst"],
        "models": ["Industriemaschinen allgemein — DST ist universal"]
    },
}


def get_embroidery_info(file_path: str) -> dict:
    """Liest Metadaten aus einer Stickdatei (Stiche, Farben, Größe)."""
    if not _PYEMB:
        return {"error": "pyembroidery nicht installiert"}
    try:
        import pyembroidery
        pattern = pyembroidery.read(file_path)
        if pattern is None:
            return {"error": "Datei konnte nicht gelesen werden"}

        # Stichzahl (nur echte Stiche, keine Kommandos)
        stitch_count = sum(1 for s in pattern.stitches if s[2] == pyembroidery.STITCH)
        jump_count   = sum(1 for s in pattern.stitches if s[2] == pyembroidery.JUMP)

        # Bounding-Box (Einheit: 0.1mm)
        xs = [s[0] for s in pattern.stitches]
        ys = [s[1] for s in pattern.stitches]
        if xs and ys:
            width_mm  = round((max(xs) - min(xs)) / 10.0, 1)
            height_mm = round((max(ys) - min(ys)) / 10.0, 1)
        else:
            width_mm = height_mm = 0.0

        # Farb-Threads
        threads = []
        for t in (pattern.threadlist or []):
            color = getattr(t, "color", 0)
            r = (color >> 16) & 0xFF
            g = (color >> 8)  & 0xFF
            b = color         & 0xFF
            threads.append({
                "hex":  f"#{r:02X}{g:02X}{b:02X}",
                "name": getattr(t, "name", "") or f"Thread {len(threads)+1}",
                "brand": getattr(t, "manufacturer", "") or "",
                "catalog": getattr(t, "catalog_number", "") or "",
            })

        return {
            "stitch_count": stitch_count,
            "jump_count":   jump_count,
            "thread_count": len(threads),
            "width_mm":     width_mm,
            "height_mm":    height_mm,
            "threads":      threads,
        }
    except Exception as e:
        return {"error": str(e)}


def get_embroidery_preview_svg(file_path: str) -> str | None:
    """Rendert eine Stickdatei als SVG-Vorschau (Fadenpfade farbig)."""
    if not _PYEMB:
        return None
    try:
        import pyembroidery
        pattern = pyembroidery.read(file_path)
        if pattern is None:
            return None

        xs = [s[0] for s in pattern.stitches]
        ys = [s[1] for s in pattern.stitches]
        if not xs:
            return None

        minx, miny = min(xs), min(ys)
        maxx, maxy = max(xs), max(ys)
        W = max(maxx - minx, 1)
        H = max(maxy - miny, 1)

        # Skalieren auf max 600px
        scale = min(600 / W, 600 / H)
        svgW = round(W * scale)
        svgH = round(H * scale)

        threads = pattern.threadlist or []
        def thread_color(idx):
            if idx < len(threads):
                c = getattr(threads[idx], "color", 0)
                r = (c >> 16) & 0xFF
                g = (c >> 8)  & 0xFF
                b = c         & 0xFF
                return f"#{r:02X}{g:02X}{b:02X}"
            # Fallback-Farbe
            fallbacks = ["#CC3333","#3355CC","#33AA55","#DDAA00","#AA33CC","#DD7700"]
            return fallbacks[idx % len(fallbacks)]

        parts = [f'<svg xmlns="http://www.w3.org/2000/svg" '
                 f'viewBox="0 0 {svgW} {svgH}" '
                 f'width="{svgW}" height="{svgH}" '
                 f'style="background:#1a1a2e;border-radius:8px">']

        thread_idx = 0
        current_path = []
        in_stitch = False

        for sx, sy, cmd in pattern.stitches:
            if cmd == pyembroidery.COLOR_CHANGE:
                # Aktiven Pfad schließen
                if current_path and len(current_path) > 1:
                    d = "M " + " L ".join(f"{round((x-minx)*scale,1)},{round((y-miny)*scale,1)}"
                                          for x, y in current_path)
                    col = thread_color(thread_idx)
                    parts.append(f'<path d="{d}" stroke="{col}" stroke-width="0.8" '
                                 f'fill="none" stroke-linecap="round" stroke-linejoin="round" opacity="0.85"/>')
                thread_idx += 1
                current_path = []
                in_stitch = False

            elif cmd == pyembroidery.STITCH:
                current_path.append((sx, sy))
                in_stitch = True

            elif cmd == pyembroidery.JUMP:
                if current_path and len(current_path) > 1:
                    d = "M " + " L ".join(f"{round((x-minx)*scale,1)},{round((y-miny)*scale,1)}"
                                          for x, y in current_path)
                    col = thread_color(thread_idx)
                    parts.append(f'<path d="{d}" stroke="{col}" stroke-width="0.8" '
                                 f'fill="none" stroke-linecap="round" stroke-linejoin="round" opacity="0.85"/>')
                current_path = [(sx, sy)]
                in_stitch = False

        # Letzten Pfad
        if current_path and len(current_path) > 1:
            d = "M " + " L ".join(f"{round((x-minx)*scale,1)},{round((y-miny)*scale,1)}"
                                   for x, y in current_path)
            col = thread_color(thread_idx)
            parts.append(f'<path d="{d}" stroke="{col}" stroke-width="0.8" '
                         f'fill="none" stroke-linecap="round" stroke-linejoin="round" opacity="0.85"/>')

        parts.append("</svg>")
        return "\n".join(parts)
    except Exception:
        return None
